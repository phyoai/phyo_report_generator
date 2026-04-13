"""
Campaign Report Generator - Flask Backend
Analyzes campaign screenshots using Anthropic Claude and populates PowerPoint template
"""

from flask import Flask, request, jsonify, send_file, Response
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.dml.color import RGBColor
from anthropic import Anthropic
import os
import sys
import base64
import io
import json
import mimetypes
import re
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor
import copy
import time
import subprocess
from dotenv import load_dotenv
import pytesseract
from PIL import Image, ImageFilter, ImageOps
import shutil

try:
    import easyocr
    HAS_EASYOCR = True
except Exception:
    easyocr = None
    HAS_EASYOCR = False

# Configure Tesseract path for Windows
import os
if os.name == 'nt':  # Windows
    pytesseract.pytesseract.pytesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

load_dotenv(override=True)  # Load environment variables from .env file (override=True forces reload)

app = Flask(__name__)
CORS(app)  # Enable CORS for React frontend

# Free tools for video/social metrics (no API keys needed)
try:
    import yt_dlp
    HAS_YT_DLP = True
except:
    HAS_YT_DLP = False

# Initialize client - supports both Anthropic and OpenRouter
if os.environ.get("OPENROUTER_API_KEY"):
    # Use OpenRouter with OpenAI-compatible format
    from openai import OpenAI as OpenAIClient
    client = OpenAIClient(
        api_key=os.environ.get("OPENROUTER_API_KEY"),
        base_url="https://openrouter.ai/api/v1",
        default_headers={
            "HTTP-Referer": "http://localhost:5000",
            "X-Title": "Campaign Report Generator"
        }
    )
    USE_OPENROUTER = True
else:
    # Use direct Anthropic API
    client = Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))
    USE_OPENROUTER = False

# Configuration
# Get the directory where this script is located
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(SCRIPT_DIR, "Campaign_Report_PyroMedia.pptx")
OUTPUT_DIR = os.path.join(SCRIPT_DIR, "generated_reports")
PREVIEW_DIR = os.path.join(OUTPUT_DIR, "template_preview")
REPORT_PREVIEW_ROOT = os.path.join(OUTPUT_DIR, "report_previews")
os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(PREVIEW_DIR, exist_ok=True)
os.makedirs(REPORT_PREVIEW_ROOT, exist_ok=True)

FETCH_CACHE_TTL_SECONDS = 300
_FETCH_CACHE = {}

class InMemoryUpload(io.BytesIO):
    def __init__(self, data, filename="", content_type=""):
        super().__init__(data)
        self.filename = filename
        self.content_type = content_type


def make_in_memory_upload(upload):
    upload_bytes = upload.read()
    upload.seek(0)
    return InMemoryUpload(
        upload_bytes,
        filename=getattr(upload, "filename", "") or "",
        content_type=getattr(upload, "content_type", "") or "",
    )


def clone_in_memory_upload(upload):
    upload.seek(0)
    return InMemoryUpload(
        upload.getvalue(),
        filename=getattr(upload, "filename", "") or "",
        content_type=getattr(upload, "content_type", "") or "",
    )


def clone_uploads(images):
    return [clone_in_memory_upload(image) for image in images]


def _cache_get(cache_key):
    cached = _FETCH_CACHE.get(cache_key)
    if not cached:
        return None
    cached_at, value = cached
    if time.time() - cached_at > FETCH_CACHE_TTL_SECONDS:
        _FETCH_CACHE.pop(cache_key, None)
        return None
    return copy.deepcopy(value)


def _cache_set(cache_key, value):
    _FETCH_CACHE[cache_key] = (time.time(), copy.deepcopy(value))


def _ps_quote(value):
    return "'" + str(value).replace("'", "''") + "'"


def _export_presentation_preview_images(pptx_path, preview_dir, image_url_builder, updated_at_key="updated_at"):
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"Presentation file not found: {pptx_path}")

    os.makedirs(preview_dir, exist_ok=True)
    source_mtime = os.path.getmtime(pptx_path)
    existing_previews = [
        name
        for name in os.listdir(preview_dir)
        if re.fullmatch(r"Slide\d+\.PNG", name, re.IGNORECASE)
    ]
    preview_mtime = 0
    if existing_previews:
        preview_mtime = min(
            os.path.getmtime(os.path.join(preview_dir, name))
            for name in existing_previews
        )

    if existing_previews and preview_mtime >= source_mtime:
        slide_names = sorted(existing_previews, key=lambda item: int(re.sub(r"\D", "", item)))
        return {
            updated_at_key: datetime.fromtimestamp(source_mtime).isoformat(),
            "slides": [
                {
                    "slideNumber": index + 1,
                    "imageUrl": image_url_builder(name, source_mtime),
                }
                for index, name in enumerate(slide_names)
            ],
        }

    for name in existing_previews:
        try:
            os.remove(os.path.join(preview_dir, name))
        except PermissionError:
            print(f"[WARNING] Could not remove locked preview image: {name}")

    if os.name != 'nt' or shutil.which("powershell") is None:
        print("[WARNING] PPT preview export is unavailable in this environment; returning empty preview set")
        return {
            updated_at_key: datetime.fromtimestamp(source_mtime).isoformat(),
            "slides": [],
            "warning": "Presentation preview rendering is unavailable on this server environment.",
        }

    preview_dir_ps = _ps_quote(preview_dir)
    template_path_ps = _ps_quote(pptx_path)
    export_script = f"""
$ErrorActionPreference = 'Stop'
$powerpoint = New-Object -ComObject PowerPoint.Application
$presentation = $null
try {{
  $presentation = $powerpoint.Presentations.Open({template_path_ps}, $false, $false, $false)
  $presentation.Export({preview_dir_ps}, 'PNG')
}} finally {{
  if ($presentation -ne $null) {{
    $presentation.Close()
  }}
  $powerpoint.Quit()
}}
"""
    subprocess.run(
        ["powershell", "-NoProfile", "-Command", export_script],
        check=True,
        capture_output=True,
        text=True,
        timeout=120,
    )

    slide_images = sorted(
        [
            name
            for name in os.listdir(preview_dir)
            if re.fullmatch(r"Slide\d+\.PNG", name, re.IGNORECASE)
        ],
        key=lambda item: int(re.sub(r"\D", "", item)),
    )

    return {
        updated_at_key: datetime.fromtimestamp(source_mtime).isoformat(),
        "slides": [
            {
                "slideNumber": index + 1,
                "imageUrl": image_url_builder(name, source_mtime),
            }
            for index, name in enumerate(slide_images)
        ],
    }


def export_template_preview_images():
    """Export the local PPT template to PNG slide previews for browser viewing."""
    return _export_presentation_preview_images(
        TEMPLATE_PATH,
        PREVIEW_DIR,
        lambda name, mtime: f"/api/template-preview/image/{name}?v={int(mtime)}",
        updated_at_key="template_updated_at",
    )


def export_report_preview_images(filename):
    """Export a generated PPTX report to PNG previews for browser viewing."""
    safe_name = os.path.basename(filename)
    if safe_name != filename or not safe_name.lower().endswith(".pptx"):
        raise ValueError("Invalid report filename")

    pptx_path = os.path.join(OUTPUT_DIR, safe_name)
    report_key = re.sub(r"[^A-Za-z0-9_.-]+", "_", os.path.splitext(safe_name)[0]).strip("._") or "report"
    report_preview_dir = os.path.join(REPORT_PREVIEW_ROOT, report_key)
    return _export_presentation_preview_images(
        pptx_path,
        report_preview_dir,
        lambda name, mtime: f"/api/report-preview/image/{report_key}/{name}?v={int(mtime)}",
        updated_at_key="report_updated_at",
    )


OPENROUTER_PROMPT_TEMPLATES = {
    "master": """You are an influencer marketing analyst.

Create a professional campaign performance report using the data below.

Campaign Details:
- Campaign Name: {campaign_name}
- Influencer: {influencer}
- Budget: ${budget}
- Duration: {duration}
- Platform: {platform}

Input Metrics:
- Views: {views}
- Likes: {likes}
- Comments: {comments}
- Total Engagement: {total_engagement}
- CPV: {cpv}
- CPE: {cpe}
- Engagement Rate: {engagement_rate}

Important Rules:
- Shares and Saves are not available, do not fabricate real values
- You may estimate them logically if needed but clearly mark them as estimated
- Keep report concise but professional
- Give actionable insights

Output Format (STRICT):
1. Campaign Summary (short paragraph)
2. Key Metrics (bullet points)
3. Performance Analysis
4. Recommendations
5. Final JSON (very important)

JSON Format:
{{
  "views": number,
  "likes": number,
  "comments": number,
  "shares": "N/A or estimated",
  "saves": "N/A or estimated",
  "total_engagement": number,
  "budget": number,
  "cpv": number,
  "cpe": number,
  "engagement_rate": number
}}""",
    "fast": """Generate a short influencer campaign report.

Data:
Views: {views}
Likes: {likes}
Comments: {comments}
Budget: {budget}

Calculate:
- Total Engagement
- CPV
- CPE
- Engagement Rate

Return:
1. Summary
2. Insights
3. JSON output""",
    "advanced": """Act as a senior marketing analyst.

Analyze this influencer campaign deeply.

Campaign:
{campaign_name}
Budget: ${budget}

Metrics:
Views: {views}
Likes: {likes}
Comments: {comments}

Tasks:
1. Calculate engagement metrics
2. Evaluate performance (low / average / high)
3. Identify strengths and weaknesses
4. Suggest improvements for next campaign
5. Compare engagement efficiency vs cost

Return:
- Performance Score (out of 10)
- Insights
- Optimization strategy
- JSON metrics""",
    "json_only": """Calculate campaign metrics and return ONLY JSON.

Input:
Views: {views}
Likes: {likes}
Comments: {comments}
Budget: {budget}

Output JSON:
{{
  "views": number,
  "likes": number,
  "comments": number,
  "total_engagement": number,
  "cpv": number,
  "cpe": number,
  "engagement_rate": number
}}""",
    "multi_platform": """Create a campaign report for:

Campaign: {campaign_name}
Budget: ${budget}

Instagram Metrics:
Views: {ig_views}
Likes: {ig_likes}
Comments: {ig_comments}

YouTube Metrics:
Views: {yt_views}
Likes: {yt_likes}
Comments: {yt_comments}

Tasks:
- Compare platform performance
- Identify which platform performed better
- Provide insights
- Return combined JSON

JSON:
{{
  "instagram": {{...}},
  "youtube": {{...}},
  "best_platform": "Instagram/YouTube"
}}""",
}


def has_tesseract():
    """Return True when the Tesseract executable is available."""
    configured_path = getattr(pytesseract.pytesseract, "tesseract_cmd", "")
    return bool(
        (configured_path and os.path.exists(configured_path))
        or shutil.which("tesseract")
    )


def detect_image_media_type(image_file, image_bytes):
    """Best-effort detection of an uploaded image's MIME type."""
    content_type = getattr(image_file, "content_type", "") or ""
    if content_type.startswith("image/"):
        return content_type

    filename = getattr(image_file, "filename", "") or ""
    guessed_type, _ = mimetypes.guess_type(filename)
    if guessed_type and guessed_type.startswith("image/"):
        return guessed_type

    try:
        pil_image = Image.open(io.BytesIO(image_bytes))
        image_format = (pil_image.format or "").lower()
        format_map = {
            "jpg": "image/jpeg",
            "jpeg": "image/jpeg",
            "png": "image/png",
            "webp": "image/webp",
            "gif": "image/gif",
            "bmp": "image/bmp",
        }
        return format_map.get(image_format, "image/jpeg")
    except Exception:
        return "image/jpeg"


class SafePromptDict(dict):
    def __missing__(self, key):
        return "N/A"


def build_openrouter_prompt(template_name="master", data=None):
    """Format one of the built-in OpenRouter prompt templates."""
    template = OPENROUTER_PROMPT_TEMPLATES.get(template_name)
    if not template:
        raise ValueError(f"Unknown prompt template: {template_name}")

    payload = SafePromptDict({
        "campaign_name": "Summer Fashion Campaign",
        "influencer": "Influencer X",
        "budget": 5000,
        "duration": "June 1 to June 30",
        "platform": "Instagram Reel",
        "views": "N/A",
        "likes": "N/A",
        "comments": "N/A",
        "total_engagement": "N/A",
        "cpv": "N/A",
        "cpe": "N/A",
        "engagement_rate": "N/A",
        "ig_views": "N/A",
        "ig_likes": "N/A",
        "ig_comments": "N/A",
        "yt_views": "N/A",
        "yt_likes": "N/A",
        "yt_comments": "N/A",
    })
    if data:
        payload.update(data)
    return template.format_map(payload)


def extract_metrics_from_image_ocr(image_file):
    """Extract metrics from image using Tesseract OCR"""
    import re

    try:
        if not has_tesseract():
            print("[WARNING]  OCR skipped: Tesseract is not installed or not available in PATH")
            return {}

        print("[SEARCH] Running OCR on image...")

        # Read image file
        image_bytes = image_file.read()
        image_file.seek(0)

        # Convert to PIL Image
        img = Image.open(io.BytesIO(image_bytes))

        def _image_variants(base_image):
            variants = []
            rgb = base_image.convert("RGB")
            gray = ImageOps.grayscale(rgb)
            enlarged = gray.resize((gray.width * 2, gray.height * 2), Image.Resampling.LANCZOS)
            high_contrast = ImageOps.autocontrast(enlarged)
            threshold = high_contrast.point(lambda x: 255 if x > 160 else 0)
            sharpened = high_contrast.filter(ImageFilter.SHARPEN)
            variants.extend([rgb, gray, high_contrast, threshold, sharpened])
            return variants

        extracted_chunks = []
        for variant in _image_variants(img):
            try:
                extracted_chunks.append(
                    pytesseract.image_to_string(
                        variant,
                        config="--oem 3 --psm 6"
                    )
                )
                extracted_chunks.append(
                    pytesseract.image_to_string(
                        variant,
                        config="--oem 3 --psm 11"
                    )
                )
            except Exception:
                continue

        if HAS_EASYOCR:
            try:
                reader = easyocr.Reader(['en'], gpu=False)
                easyocr_text = reader.readtext(image_bytes, detail=0, paragraph=True)
                if easyocr_text:
                    extracted_chunks.append("\n".join(easyocr_text))
            except Exception as easyocr_error:
                print(f"[WARNING] EasyOCR failed: {easyocr_error}")

        extracted_text = "\n".join(chunk for chunk in extracted_chunks if chunk).strip()

        print(f"[NOTE] OCR extracted text (first 300 chars):\n{extracted_text[:300]}...\n")

        # Parse metrics from extracted text
        return extract_metrics_from_text(extracted_text)

    except Exception as e:
        print(f"[WARNING]  OCR failed: {str(e)}")
        return {}


def extract_metrics_from_images_ocr(images):
    merged_metrics = {}
    if not images:
        return merged_metrics

    max_workers = min(4, len(images)) or 1
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = [executor.submit(extract_metrics_from_image_ocr, image) for image in images]
        for future in futures:
            image_metrics = future.result() or {}
            for key, value in image_metrics.items():
                if value and value != "N/A" and not merged_metrics.get(key):
                    merged_metrics[key] = value
    return merged_metrics


def build_remote_media_uploads(instagram_post=None):
    uploads = []
    if not instagram_post:
        return uploads

    media_urls = []
    for key in ["postImage", "videoUrl"]:
        value = instagram_post.get(key)
        if value and value not in media_urls:
            media_urls.append(value)

    for index, media_url in enumerate(media_urls):
        try:
            import requests

            response = requests.get(
                media_url,
                timeout=30,
                headers={
                    "User-Agent": (
                        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                        "AppleWebKit/537.36 (KHTML, like Gecko) "
                        "Chrome/123.0.0.0 Safari/537.36"
                    )
                },
            )
            response.raise_for_status()
            content_type = response.headers.get("Content-Type", "") or "image/jpeg"
            if not content_type.startswith("image/"):
                continue
            uploads.append(
                InMemoryUpload(
                    response.content,
                    filename=f"instagram_media_{index + 1}.jpg",
                    content_type=content_type,
                )
            )
        except Exception as e:
            print(f"[WARNING] Could not download Instagram media for OCR fallback: {e}")

    return uploads


def fetch_youtube_metrics(url_or_text):
    """
    Fetch YouTube metrics using yt-dlp (FREE - no API key needed).
    Handles both video URLs and channel/username URLs (@handle or /channel/).
    Returns dict with keys matching the report's youtube schema.
    """
    if not HAS_YT_DLP:
        print("[WARNING]  yt-dlp not installed")
        return {}

    import yt_dlp
    import re

    # Accept either a raw URL or a blob of text containing a URL
    youtube_url = url_or_text.strip()
    if not youtube_url.startswith("http"):
        urls = re.findall(r'(https?://(?:www\.)?(?:youtube\.com|youtu\.be)[^\s]*)', url_or_text)
        if not urls:
            return {}
        youtube_url = urls[0]

    print(f"[VIDEO] Fetching YouTube metrics from: {youtube_url}")

    # Detect whether this is a channel/user URL or a direct video URL
    is_channel = any(x in youtube_url for x in [
        "/@", "/channel/", "/user/", "/c/",
        "youtube.com/@", "youtube.com/channel",
    ])

    quiet_opts = {"quiet": True, "no_warnings": True}

    try:
        if is_channel:
            # --- Channel URL: fetch the most recent upload and use its metrics ---
            channel_opts = {
                **quiet_opts,
                "extract_flat": "in_playlist",   # list videos without downloading
                "playlistend": 1,                 # only the latest video
            }
            with yt_dlp.YoutubeDL(channel_opts) as ydl:
                channel_info = ydl.extract_info(youtube_url, download=False)

            # Get the first video URL from the channel playlist
            entries = channel_info.get("entries") or []
            if not entries:
                print("[WARNING]  No videos found on channel")
                return {}

            first_entry = entries[0]
            video_id  = first_entry.get("id") or first_entry.get("url", "")
            video_url = f"https://www.youtube.com/watch?v={video_id}" if len(video_id) == 11 else first_entry.get("url", "")
            creator_name = channel_info.get("uploader") or channel_info.get("channel") or ""
            print(f"[VIDEO] Latest video: {video_url}")
        else:
            video_url    = youtube_url
            creator_name = ""

        # --- Fetch full video metadata ---
        with yt_dlp.YoutubeDL(quiet_opts) as ydl:
            info = ydl.extract_info(video_url, download=False)

        metrics = {}

        if info.get("view_count"):
            metrics["views"] = str(info["view_count"])
        if info.get("like_count"):
            metrics["likes"] = str(info["like_count"])
        if info.get("comment_count"):
            metrics["comments"] = str(info["comment_count"])

        # Duration in seconds → format as "Xm Ys" (closest public proxy for watch time)
        duration = info.get("duration")
        if duration:
            mins, secs = divmod(int(duration), 60)
            metrics["watchTime"] = f"{mins}m {secs}s"

        # Top-level metadata (used outside youtube dict)
        metrics["_title"]        = info.get("title", "")
        metrics["_creator_name"] = creator_name or info.get("uploader") or info.get("channel") or ""

        print(f"[OK] YouTube metrics: {metrics}")
        return metrics

    except Exception as e:
        print(f"[WARNING]  YouTube fetch error: {str(e)}")
        return {}


def extract_metrics_from_text(text):
    """Extract campaign metrics from user's text input - AGGRESSIVE extraction"""
    import re

    data = {
        "clicks": "N/A",
        "views": "N/A",
        "likes": "N/A",
        "shares": "N/A",
        "saves": "N/A",
        "total_engagement": "N/A",
        "budget": "N/A",
        "budget_currency": "",
        "cpv": "N/A",
        "cpe": "N/A",
        "cpc": "N/A",
        "cpc_goal": "",
        "cpv_goal": "",
        "cpc_calculation": "",
        "cpv_calculation": "",
        "engagement_rate": "N/A"
    }

    text_lower = text.lower()

    # Extract Views (multiple patterns)
    views_patterns = [
        r'no\.?\s+of\s+views?[\s:]*([\d,\.]+[kmb]?)',
        r'views?[\s:]*[\$]?([\d,\.]+[kmb]?)',
        r'(\d+[,\d]*)\s+views?',
        r'view count[\s:]*(\d+[,\d]*)',
    ]
    for pattern in views_patterns:
        views_match = re.search(pattern, text_lower)
        if views_match:
            data["views"] = views_match.group(1).replace(',', '').replace('k', '000').replace('m', '000000')
            break

    clicks_patterns = [
        r'clicks?[\s:]*([\d,\.]+[kmb]?)',
        r'total\s+clicks?[\s:]*([\d,\.]+[kmb]?)',
        r'link\s+clicks?[\s:]*([\d,\.]+[kmb]?)',
    ]
    for pattern in clicks_patterns:
        clicks_match = re.search(pattern, text_lower)
        if clicks_match:
            data["clicks"] = clicks_match.group(1).replace(',', '').replace('k', '000').replace('m', '000000')
            break

    # Extract Likes (multiple patterns)
    likes_patterns = [
        r'likes?[\s:]*[\$]?([\d,\.]+[kmb]?)',
        r'(\d+[,\d]*)\s+likes?',
        r'like count[\s:]*(\d+[,\d]*)',
    ]
    for pattern in likes_patterns:
        likes_match = re.search(pattern, text_lower)
        if likes_match:
            data["likes"] = likes_match.group(1).replace(',', '').replace('k', '000').replace('m', '000000')
            break

    # Extract Shares
    shares_patterns = [
        r'shares?[\s:]*[\$]?([\d,\.]+[kmb]?)',
        r'(\d+[,\d]*)\s+shares?',
    ]
    for pattern in shares_patterns:
        shares_match = re.search(pattern, text_lower)
        if shares_match:
            data["shares"] = shares_match.group(1).replace(',', '').replace('k', '000').replace('m', '000000')
            break

    # Extract Saves
    saves_patterns = [
        r'saves?[\s:]*[\$]?([\d,\.]+[kmb]?)',
        r'(\d+[,\d]*)\s+saves?',
    ]
    for pattern in saves_patterns:
        saves_match = re.search(pattern, text_lower)
        if saves_match:
            data["saves"] = saves_match.group(1).replace(',', '').replace('k', '000').replace('m', '000000')
            break

    # Extract Budget (multiple patterns)
    budget_patterns = [
        (r'budget[\s:]*([\d,]+\.?\d*)\s*inr', 'INR'),
        (r'budget[\s:]*inr\s*([\d,]+\.?\d*)', 'INR'),
        (r'budget[\s:]*rs\.?\s*([\d,]+\.?\d*)', 'INR'),
        (r'budget[\s:]*₹\s*([\d,]+\.?\d*)', 'INR'),
        (r'budget[\s:]*[\$]([\d,]+\.?\d*)', 'USD'),
        (r'\$([\d,]+)\s+budget', 'USD'),
        (r'budget.*?([\d,]+)', ''),
    ]
    for pattern, currency in budget_patterns:
        budget_match = re.search(pattern, text_lower)
        if budget_match:
            data["budget"] = budget_match.group(1).replace(',', '')
            data["budget_currency"] = currency
            break

    # Extract CPV (Cost Per View)
    cpv_match = re.search(r'(?:cpv|cost\s+per\s+view)[\s:]*?(?:rs\.?|inr|\$)?\s*([\d,]+(?:\.\d+)?)', text_lower)
    if cpv_match:
        data["cpv"] = cpv_match.group(1).replace(',', '')

    # Extract CPC (Cost Per Click)
    cpc_match = re.search(r'(?:cpc|cost\s+per\s+click)[\s:]*?(?:rs\.?|inr|\$)?\s*([\d,]+(?:\.\d+)?)', text_lower)
    if cpc_match:
        data["cpc"] = cpc_match.group(1).replace(',', '')

    # Extract CPE (Cost Per Engagement)
    cpe_match = re.search(r'(?:cpe|cost\s+per\s+engagement)[\s:]*?(?:rs\.?|inr|\$)?\s*([\d,]+(?:\.\d+)?)', text_lower)
    if cpe_match:
        data["cpe"] = cpe_match.group(1).replace(',', '')

    # Extract explicit Total Engagement if provided.
    total_engagement_match = re.search(r'total\s+engagement[\s:]*([\d,\.]+[kmb]?)', text_lower)
    if total_engagement_match:
        data["total_engagement"] = total_engagement_match.group(1).replace(',', '').replace('k', '000').replace('m', '000000')

    # Capture non-numeric goal/intent lines for CPC / CPV if present.
    for metric_key, goal_key in [("cpc", "cpc_goal"), ("cpv", "cpv_goal")]:
        goal_match = re.search(rf'(?im)^\s*{metric_key}\s*:\s*([^\n]*[A-Za-z][^\n]*)$', text)
        if goal_match:
            data[goal_key] = goal_match.group(1).strip()

    calc_section_match = re.search(r'(?is)calculation\s*:\s*(.+)', text)
    if calc_section_match:
        calc_section = calc_section_match.group(1)
        cpc_calc_match = re.search(r'(?im)^\s*cpc\s*:\s*([^\n]+)$', calc_section)
        cpv_calc_match = re.search(r'(?im)^\s*cpv\s*:\s*([^\n]+)$', calc_section)
        if cpc_calc_match:
            data["cpc_calculation"] = cpc_calc_match.group(1).strip()
        if cpv_calc_match:
            data["cpv_calculation"] = cpv_calc_match.group(1).strip()

    # Extract Engagement Rate
    eng_rate_patterns = [
        r'engagement\s+rate[\s:]*(\d+\.?\d*)%?',
        r'engagement rate.*?(\d+\.?\d*)%',
    ]
    for pattern in eng_rate_patterns:
        eng_rate_match = re.search(pattern, text_lower)
        if eng_rate_match:
            data["engagement_rate"] = eng_rate_match.group(1)
            break

    # Calculate total engagement if not provided
    if data["total_engagement"] == "N/A":
        engagement_sum = 0
        count = 0
        for metric in ["likes", "shares", "saves"]:
            if data[metric] != "N/A":
                try:
                    engagement_sum += int(float(data[metric]))
                    count += 1
                except:
                    pass
        if count > 0:
            data["total_engagement"] = str(int(engagement_sum))

    print(f"[OK] Text metrics extracted: {data}")
    return data


def extract_prompt_context(prompt):
    """Extract campaign identity fields from the user's text prompt."""
    import re

    context = {
        "campaignName": "",
        "brand": "",
        "creator": "",
        "startDate": "",
        "endDate": "",
        "deliverables": "",
        "financial": {},
    }

    if not prompt:
        return context

    def first(patterns, text, flags=re.IGNORECASE):
        for pattern in patterns:
            match = re.search(pattern, text, flags)
            if match:
                return match.group(1).strip()
        return ""

    campaign_name = first([
        r'report\s+for\s+([^\n,\.]+?campaign)(?:\s+with\b|\s*,|\s*$)',
        r'campaign\s+name\s*[:\-]\s*([^\n,\.]+)',
        r'\bcampaign\s*:\s*([^\n,\.]+)',
        r'for\s+([^\n,\.]+?campaign)(?:\s+with\b|\s*,|\s*$)',
    ], prompt)
    if campaign_name:
        context["campaignName"] = campaign_name[:80]

    brand = first([r'\bbrand\s*[:\-]\s*([^\n,\.]+)'], prompt)
    if brand:
        context["brand"] = brand[:40]
    else:
        repeated_name_matches = re.findall(r'(?im)^\s*([A-Z][A-Za-z0-9&.+-]{2,})\s*$', prompt)
        repeated_names = []
        for name in repeated_name_matches:
            if repeated_name_matches.count(name) >= 2 and name.lower() not in {"cpc", "cpv", "budget", "campaign"}:
                repeated_names.append(name)
        if repeated_names:
            context["brand"] = repeated_names[0][:40]

    creator = first([
        r'(?:influencer|creator)\s*[:\-]\s*([^\n,\.]{2,})',
        r'with\s+((?:influencer|creator)\s+[^\n,\.]+?)(?:\s*,|\s+budget|\s+ran|\s+from|$)',
        r'with\s+(@[A-Za-z0-9_.]+)',
    ], prompt)
    if creator and len(creator.strip()) > 1:
        context["creator"] = creator[:60]

    range_match = re.search(
        r'ran\s+([A-Za-z]+)\s+(\d+)\s*(?:-|to|through|until)\s*(?:([A-Za-z]+)\s+)?(\d+)',
        prompt,
        re.IGNORECASE,
    )
    if range_match:
        month_start = range_match.group(1)
        day_start = range_match.group(2)
        month_end = range_match.group(3) or month_start
        day_end = range_match.group(4)
        context["startDate"] = f"{month_start} {day_start}"
        context["endDate"] = f"{month_end} {day_end}"
    else:
        start_date = first([r'(?:ran|from|start(?:ed)?)\s+([A-Za-z]+\s+\d+)'], prompt)
        end_date = first([r'(?:to|through|until|end(?:ed)?)\s+([A-Za-z]+\s+\d+)'], prompt)
        if start_date:
            context["startDate"] = start_date
        if end_date:
            context["endDate"] = end_date

    budget = first([
        r'budget\s*\$?([\d,]+(?:\.\d+)?)',
        r'\$([\d,]+(?:\.\d+)?)\s+budget',
    ], prompt)
    if budget:
        context["financial"]["totalBudget"] = budget.replace(",", "")
    if re.search(r'\b(inr|rs\.?|₹)\b', prompt, re.IGNORECASE):
        context["financial"]["budgetCurrency"] = "INR"
    elif "$" in prompt:
        context["financial"]["budgetCurrency"] = "USD"

    cpv = first([r'\bcpv\s*[:\-]?\s*\$?([\d.]+)'], prompt)
    if cpv:
        context["financial"]["cpv"] = cpv

    cpe = first([r'\bcpe\s*[:\-]?\s*\$?([\d.]+)'], prompt)
    if cpe:
        context["financial"]["cpe"] = cpe

    cpc = first([r'\bcpc\s*[:\-]?\s*\$?([\d.]+)'], prompt)
    if cpc:
        context["financial"]["cpc"] = cpc

    for metric_key, fin_key in [("cpc", "cpcGoal"), ("cpv", "cpvGoal")]:
        goal_match = re.search(rf'(?im)^\s*{metric_key}\s*:\s*([^\n]*[A-Za-z][^\n]*)$', prompt)
        if goal_match:
            context["financial"][fin_key] = goal_match.group(1).strip()

    calc_section_match = re.search(r'(?is)calculation\s*:\s*(.+)', prompt)
    if calc_section_match:
        calc_section = calc_section_match.group(1)
        for metric_key, fin_key in [("cpc", "cpcCalculation"), ("cpv", "cpvCalculation")]:
            calc_match = re.search(rf'(?im)^\s*{metric_key}\s*:\s*([^\n]+)$', calc_section)
            if calc_match:
                context["financial"][fin_key] = calc_match.group(1).strip()

    mentioned_platforms = []
    if re.search(r'\binstagram\b', prompt, re.IGNORECASE):
        mentioned_platforms.append("Instagram")
    if re.search(r'\byoutube\b', prompt, re.IGNORECASE):
        mentioned_platforms.append("YouTube")
    if mentioned_platforms:
        context["deliverables"] = " + ".join(mentioned_platforms)

    return context


def apply_prompt_context(extracted_data, prompt_context, force_identity=False):
    """Apply prompt-derived campaign context onto extracted data."""
    if not isinstance(extracted_data, dict):
        extracted_data = create_default_data()

    if not prompt_context:
        return extracted_data

    for key in ["campaignName", "creator", "startDate", "endDate"]:
        value = prompt_context.get(key)
        if value and (force_identity or not extracted_data.get(key)):
            extracted_data[key] = value

    # Brand should only come from the prompt when the prompt explicitly provides it.
    brand = prompt_context.get("brand")
    if brand and (force_identity or not extracted_data.get("brand") or extracted_data.get("brand") == "Unknown Brand"):
        extracted_data["brand"] = brand

    deliverables = prompt_context.get("deliverables")
    if deliverables and not extracted_data.get("deliverables"):
        extracted_data["deliverables"] = deliverables

    for fin_key in ["totalBudget", "budgetCurrency", "cpv", "cpe", "cpc", "cpcGoal", "cpvGoal", "cpcCalculation", "cpvCalculation", "roi"]:
        fin_value = prompt_context.get("financial", {}).get(fin_key)
        if fin_value and not extracted_data.get("financial", {}).get(fin_key):
            extracted_data.setdefault("financial", {})[fin_key] = fin_value

    return extracted_data


def has_meaningful_metrics(extracted_data):
    """Return True if any Instagram or YouTube metric fields contain values."""
    if not isinstance(extracted_data, dict):
        return False

    def _has_value(value):
        if value is None:
            return False
        if isinstance(value, str):
            return value.strip() not in ("", "0", "0.0", "0%")
        return bool(value)

    metric_fields = {
        "instagram": ["views", "likes", "comments", "shares", "saves", "reach", "impressions", "engagementRate"],
        "youtube": ["views", "likes", "comments", "shares", "watchTime", "ctr"],
    }
    for section, keys in metric_fields.items():
        values = extracted_data.get(section, {})
        if any(_has_value(values.get(key)) for key in keys):
            return True
    return False


def create_default_data():
    """Return PyroMedia style data structure (matches React component)"""
    return {
        "campaignName": "Campaign Report",
        "brand": "Unknown Brand",
        "agency": "PyroMedia",
        "agencyContact": "",
        "agencyWebsite": "",
        "creator": "",
        "startDate": "",
        "endDate": "",
        "deliverables": "",
        "postImage": "",
        "videoUrl": "",
        "instagram": {
            "views": "",
            "likes": "",
            "comments": "",
            "shares": "",
            "saves": "",
            "reach": "",
            "impressions": "",
            "engagementRate": ""
        },
        "youtube": {
            "views": "",
            "likes": "",
            "comments": "",
            "shares": "",
            "watchTime": "",
            "ctr": ""
        },
        "financial": {
            "totalBudget": "",
            "budgetCurrency": "",
            "cpv": "",
            "cpe": "",
            "cpc": "",
            "cpcGoal": "",
            "cpvGoal": "",
            "cpcCalculation": "",
            "cpvCalculation": "",
            "roi": ""
        },
        "performance": {
            "totalViews": "",
            "totalLikes": "",
            "totalComments": "",
            "totalShares": "",
            "totalSaves": "",
            "totalReach": "",
            "totalClicks": "",
            "totalInteractions": "",
            "totalEngagement": "",
            "keyLearnings": "Campaign data extraction complete."
        },
        "creators": []
    }


def validate_and_clean_data(data):
    """Ensure extracted data has all required fields and proper structure"""
    if not isinstance(data, dict):
        return create_default_data()

    defaults = create_default_data()
    cleaned_data = data.copy()

    # Validate and initialize nested objects
    for key in ["instagram", "youtube", "financial", "performance"]:
        if key not in cleaned_data or not isinstance(cleaned_data[key], dict):
            cleaned_data[key] = defaults[key].copy()
        else:
            # Ensure all sub-keys from defaults exist and convert None values
            for sub_key in defaults[key]:
                if sub_key not in cleaned_data[key]:
                    cleaned_data[key][sub_key] = defaults[key][sub_key]
                elif cleaned_data[key][sub_key] is None:
                    cleaned_data[key][sub_key] = ""

    # Ensure creators is a list
    if not isinstance(cleaned_data.get("creators"), list):
        cleaned_data["creators"] = []

    # Ensure top-level fields
    for key in ["campaignName", "brand", "agency", "creator", "startDate", "endDate", "deliverables", "agencyContact", "agencyWebsite"]:
        if key not in cleaned_data:
            cleaned_data[key] = defaults.get(key, "")
        # Convert None to empty string
        if cleaned_data[key] is None:
            cleaned_data[key] = ""

    cleaned_brand = _select_unique_brand_name([cleaned_data.get("brand")])
    if cleaned_brand:
        cleaned_data["brand"] = cleaned_brand

    return cleaned_data


def calculate_metrics(data):
    """Calculate CPV and CPE if budget and views/engagement are available"""
    
    def safe_float(value):
        """Convert string to float, handling various formats"""
        if value == "N/A" or value is None or value == "":
            return None
        try:
            # Remove commas and convert to float
            return float(str(value).replace(',', '').strip())
        except (ValueError, AttributeError):
            return None
    
    def format_number(value):
        """Format number with 2 decimal places"""
        if value is None:
            return "N/A"
        return f"{value:.2f}"
    
    def format_integer(value):
        """Format number as integer"""
        if value is None:
            return "N/A"
        return str(int(value))
    
    # Calculate for overall campaign
    overall = data.get('overall_campaign', {})
    budget = safe_float(overall.get('budget'))
    views = safe_float(overall.get('views'))
    likes = safe_float(overall.get('likes'))
    shares = safe_float(overall.get('shares'))
    saves = safe_float(overall.get('saves'))
    engagement = safe_float(overall.get('total_engagement'))
    
    # Calculate total engagement if it's N/A but we have individual metrics
    if engagement is None:
        engagement_sum = 0
        count = 0
        for metric in [likes, shares, saves]:
            if metric is not None:
                engagement_sum += metric
                count += 1
        
        if count > 0:
            engagement = engagement_sum
            overall['total_engagement'] = format_integer(engagement)
    
    # Calculate CPV
    if budget and views and views > 0:
        cpv = budget / views
        overall['cpv'] = format_number(cpv)
    
    # Calculate CPE
    if budget and engagement and engagement > 0:
        cpe = budget / engagement
        overall['cpe'] = format_number(cpe)
    
    # Calculate for creator data
    creator = data.get('creator_data', {})
    
    # If creator data is N/A or not provided, copy from overall campaign
    creator_views = safe_float(creator.get('views'))
    creator_likes = safe_float(creator.get('likes'))
    creator_shares = safe_float(creator.get('shares'))
    creator_saves = safe_float(creator.get('saves'))
    creator_engagement = safe_float(creator.get('total_engagement'))
    creator_budget = safe_float(creator.get('budget'))
    
    # Copy overall data to creator if creator data is N/A
    if creator_views is None and views is not None:
        creator['views'] = overall.get('views')
        creator_views = views
    
    if creator_likes is None and likes is not None:
        creator['likes'] = overall.get('likes')
        creator_likes = likes
    
    if creator_shares is None and shares is not None:
        creator['shares'] = overall.get('shares')
        creator_shares = shares
    
    if creator_saves is None and saves is not None:
        creator['saves'] = overall.get('saves')
        creator_saves = saves
    
    if creator_budget is None and budget is not None:
        creator['budget'] = overall.get('budget')
        creator_budget = budget
    
    # Calculate creator total engagement if it's N/A
    if creator_engagement is None:
        engagement_sum = 0
        count = 0
        for metric in [creator_likes, creator_shares, creator_saves]:
            if metric is not None:
                engagement_sum += metric
                count += 1
        
        if count > 0:
            creator_engagement = engagement_sum
            creator['total_engagement'] = format_integer(creator_engagement)
    
    # Calculate creator CPV
    if creator_budget and creator_views and creator_views > 0:
        creator_cpv = creator_budget / creator_views
        creator['cpv'] = format_number(creator_cpv)
    
    # Calculate creator CPE
    if creator_budget and creator_engagement and creator_engagement > 0:
        creator_cpe = creator_budget / creator_engagement
        creator['cpe'] = format_number(creator_cpe)
    
    return data


def analyze_images_with_gpt4(images, user_prompt):
    """Analyze campaign screenshots using Claude Vision or OpenRouter"""
    api_name = "OpenRouter" if USE_OPENROUTER else "Claude"
    print(f"[IMAGE] Analyzing {len(images)} images with {api_name} Vision...")

    # Prepare images in the correct format
    image_contents = []
    for image in images:
        image_bytes = image.read()
        image.seek(0)  # Reset file pointer for later use
        image_b64 = base64.b64encode(image_bytes).decode('utf-8')
        media_type = detect_image_media_type(image, image_bytes)

        if USE_OPENROUTER:
            # OpenAI-compatible format for OpenRouter
            image_contents.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:{media_type};base64,{image_b64}"
                }
            })
        else:
            # Anthropic format
            image_contents.append({
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": media_type,
                    "data": image_b64
                }
            })
    
    # Build a clear extraction prompt
    analysis_prompt = f"""You are a precise data extraction assistant. Your job is to read EXACT numbers visible in these screenshots.

Campaign context from user: "{user_prompt}"

CRITICAL RULES:
1. Read the EXACT numbers shown on screen — do NOT estimate or calculate averages
2. These are screenshots of a specific Instagram post/Reel or YouTube video — extract the EXACT metrics shown
3. Numbers may appear as "3,245" or "3.2K" or "3245" — convert all to plain integers (e.g. "3245")
4. Look for numbers next to labels: Views, Plays, Likes, Comments, Shares, Saves, Reach, Impressions, Followers, Engagement Rate, Watch Time, CTR
5. If a number is NOT visible in any image, use "" — never guess

WHAT TO LOOK FOR IN INSTAGRAM SCREENSHOTS:
- Post/Reel plays or views count (shown as "X plays" or "X views")
- Likes count (heart icon number)
- Comments count (speech bubble number)
- Shares count (arrow icon number)
- Saves count (bookmark icon number)
- Reach (unique accounts reached)
- Impressions (total times shown)
- Engagement rate (shown as X%)
- Account followers count
- Username / handle (@username)

WHAT TO LOOK FOR IN YOUTUBE SCREENSHOTS:
- View count
- Like count
- Comment count
- Watch time (hours/minutes)
- CTR (click-through rate %)
- Subscribers count

Return ONLY valid JSON in exactly this format — no markdown, no extra text:
{{
  "campaignName": "exact campaign name from image or user context",
  "brand": "brand name visible in image",
  "agency": "agency name if visible",
  "creator": "exact @username or creator name from image",
  "startDate": "date if visible",
  "endDate": "date if visible",
  "deliverables": "type of content eg Instagram Reel, YouTube Video",
  "instagram": {{
    "views": "EXACT number from image",
    "likes": "EXACT number from image",
    "comments": "EXACT number from image",
    "shares": "EXACT number from image",
    "saves": "EXACT number from image",
    "reach": "EXACT number from image",
    "impressions": "EXACT number from image",
    "engagementRate": "EXACT percentage from image eg 1.82%"
  }},
  "youtube": {{
    "views": "EXACT number from image",
    "likes": "EXACT number from image",
    "comments": "EXACT number from image",
    "shares": "EXACT number from image",
    "watchTime": "EXACT watch time from image",
    "ctr": "EXACT CTR percentage from image"
  }},
  "financial": {{
    "totalBudget": "from user text if provided",
    "cpv": "if visible",
    "cpe": "if visible",
    "cpc": "if visible",
    "cpcGoal": "goal/intent text for CPC if visible",
    "cpvGoal": "goal/intent text for CPV if visible",
    "cpcCalculation": "CPC calculation text if visible",
    "cpvCalculation": "CPV calculation text if visible",
    "roi": "if visible"
  }},
  "performance": {{
    "totalViews": "sum of all views across all posts/images",
    "totalLikes": "sum of all likes",
    "totalComments": "sum of all comments",
    "totalShares": "sum of all shares",
    "totalSaves": "sum of all saves",
    "totalReach": "total reach",
    "totalInteractions": "likes + comments + shares + saves",
    "totalEngagement": "engagement rate or total engagement number",
    "keyLearnings": "2 sentences summarising what the campaign achieved based on the numbers"
  }},
  "creators": [
    {{
      "name": "creator name or @username",
      "platform": "Instagram or YouTube",
      "views": "EXACT number",
      "likes": "EXACT number",
      "comments": "EXACT number",
      "shares": "EXACT number",
      "saves": "EXACT number",
      "reach": "EXACT number",
      "engagementRate": "EXACT percentage",
      "watchTime": "if YouTube",
      "interactions": "likes + comments + shares + saves"
    }}
  ]
}}"""

    # Create messages for Claude API
    messages = [
        {
            "role": "user",
            "content": [
                {"type": "text", "text": analysis_prompt},
                *image_contents
            ]
        }
    ]

    try:
        # Call Vision API (Claude or OpenRouter)
        model_id = os.environ.get("MODEL_ID", "claude-3-5-sonnet-20241022")

        # Fallback model list in case primary model fails
        FALLBACK_MODELS = [
            "meta-llama/llama-3.2-11b-vision-instruct:free",
            "qwen/qwen-2-vl-7b-instruct:free",
            "nvidia/llama-3.1-nemotron-nano-8b-v1:free",
            "nvidia/nemotron-nano-12b-v2-vl",
        ]

        if USE_OPENROUTER:
            # Try primary model, then fallbacks if it fails
            models_to_try = [model_id] + [m for m in FALLBACK_MODELS if m != model_id]
            response = None
            used_model = None

            for try_model in models_to_try:
                try:
                    print(f"[NOTE] Trying model: {try_model}")
                    response = client.chat.completions.create(
                        model=try_model,
                        messages=[
                            {
                                "role": "system",
                                "content": "You are a data extraction assistant. Read text and numbers from images. Return ONLY valid JSON, no markdown."
                            },
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": analysis_prompt},
                                    *image_contents
                                ]
                            }
                        ],
                        max_tokens=4096,
                        temperature=0.1
                    )
                    used_model = try_model
                    print(f"[OK] Using model: {used_model}")
                    break
                except Exception as model_err:
                    print(f"[WARNING] Model {try_model} failed: {str(model_err)[:100]}")
                    continue

            if response is None:
                print("[ERROR] All models failed")
                return create_default_data()

            # Handle both regular responses and reasoning-based responses
            message = response.choices[0].message
            response_text = message.content

            # If content is None but reasoning exists, try to extract JSON from reasoning
            if not response_text and hasattr(message, 'reasoning') and message.reasoning:
                print(f"[NOTE] Extracting JSON from reasoning field...")
                response_text = message.reasoning
                if isinstance(response_text, list) and len(response_text) > 0:
                    response_text = response_text[0].get('text', '') if isinstance(response_text[0], dict) else str(response_text[0])
        else:
            # Anthropic API call
            response = client.messages.create(
                model=model_id,
                max_tokens=2000,
                messages=messages,
                system="You are a data extraction assistant specialized in reading text and numbers from marketing analytics screenshots. You extract only visible text and metrics, never identifying people."
            )
            response_text = response.content[0].text

        # Verify response is not None
        if not response_text:
            print(f"[ERROR] {api_name} returned empty response")
            print(f"Full response object: {response}")
            return create_default_data()

        print(f"[OK] {api_name} Response received ({len(response_text)} chars total)")
        print(f"[DEBUG] Full response:\n{response_text}\n[/DEBUG]")
        
        # Clean up response (remove markdown code blocks if present)
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0]
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0]

        # Extract JSON from response - find first { to last }
        response_text = response_text.strip()
        start = response_text.find('{')

        if start == -1:
            print("[ERROR] No JSON found in response")
            print(f"Full response (first 500 chars): {response_text[:500]}")
            return create_default_data()

        # Try to find the closing brace, accounting for nested structures
        brace_count = 0
        end = start
        for i in range(start, len(response_text)):
            if response_text[i] == '{':
                brace_count += 1
            elif response_text[i] == '}':
                brace_count -= 1
                if brace_count == 0:
                    end = i + 1
                    break

        # Handle truncated JSON (no matching closing brace)
        if end == start:
            print("[WARNING] JSON appears truncated - attempting repair...")
            json_str = response_text[start:]
            open_braces = json_str.count('{') - json_str.count('}')
            open_brackets = json_str.count('[') - json_str.count(']')
            json_str = json_str.rstrip(',').rstrip('"').rstrip(' ')
            # Close any open string by checking if last char is inside a string
            if open_brackets > 0:
                json_str += '"' + ']' * open_brackets
            json_str += '}' * open_braces
            print(f"[NOTE] Added {open_braces} closing braces, {open_brackets} brackets")
        else:
            json_str = response_text[start:end]

        print(f"[DATA] Extracted JSON ({len(json_str)} chars)")

        try:
            extracted_data = json.loads(json_str)
        except json.JSONDecodeError as e:
            print(f"[ERROR] JSON Parse Error (attempt 1): {str(e)}")
            # Try to fix common JSON issues (escape sequences, newlines)
            json_str_fixed = json_str.replace('\\n', ' ').replace('\n', ' ')
            try:
                extracted_data = json.loads(json_str_fixed)
                print("[OK] Fixed JSON with escape sequence cleanup")
            except:
                # Last resort: extract key fields using regex
                print("[NOTE] Attempting regex field extraction from partial JSON...")
                import re
                extracted_data = create_default_data()
                # Extract visible string fields from partial JSON
                for field in ['campaignName', 'brand', 'agency', 'creator', 'startDate', 'endDate', 'deliverables']:
                    m = re.search(rf'"{field}"\s*:\s*"([^"]*)"', json_str)
                    if m and m.group(1):
                        extracted_data[field] = m.group(1)
                # Extract instagram metrics
                for field in ['views', 'likes', 'comments', 'shares', 'saves', 'reach', 'impressions', 'engagementRate']:
                    m = re.search(rf'(?:instagram[^}}]+)?"{field}"\s*:\s*"([^"]*)"', json_str)
                    if m and m.group(1):
                        extracted_data['instagram'][field] = m.group(1)
                print(f"[NOTE] Regex extraction found: campaignName={extracted_data.get('campaignName')}, brand={extracted_data.get('brand')}")
                if extracted_data.get('brand') == 'Unknown Brand':
                    return create_default_data()
        
        # Strip "0" string values from Vision AI output — these mean "not found", not actual zeros
        def _strip_zeros(d):
            if not isinstance(d, dict):
                return
            for k, v in list(d.items()):
                if isinstance(v, str) and v.strip() in ("0", "0.0", "0%"):
                    d[k] = ""
                elif isinstance(v, dict):
                    _strip_zeros(v)
                elif isinstance(v, list):
                    for item in v:
                        _strip_zeros(item)

        _strip_zeros(extracted_data)
        print("[NOTE] Stripped '0' placeholder values from Vision AI output")

        # Fill missing top-level fields from user prompt text
        import re as _re

        def _first(patterns, text, flags=_re.IGNORECASE):
            for p in patterns:
                m = _re.search(p, text, flags)
                if m:
                    return m.group(1).strip()
            return ""

        if not extracted_data.get("campaignName") or extracted_data["campaignName"] in ("", "Campaign Report"):
            val = _first([
                # explicit label: "campaign name: X"
                r'campaign\s+name\s*[:\-]\s*([^\n,\.]+)',
                # explicit label: "campaign: X"
                r'\bcampaign\s*:\s*([^\n,\.]+)',
                # "for X campaign" — capture the descriptive words BEFORE "campaign"
                r'for\s+((?:\w+\s+){1,5}?campaign)\b',
                # "for X campaign with ..." — same but with trailing context
                r'report\s+for\s+([^\n,\.]+?campaign[^\n,\.]*?)(?:\s+with\b|\s*,|\s*$)',
            ], user_prompt)
            if val:
                extracted_data["campaignName"] = val.strip()[:60]

        if not extracted_data.get("brand") or extracted_data["brand"] in ("", "Unknown Brand"):
            val = _first([r'\bbrand[:\s]+([^\n,\.]+)'], user_prompt)
            if val:
                extracted_data["brand"] = val[:40]

        if not extracted_data.get("creator"):
            val = _first([
                # explicit labels
                r'influencer\s*[:\-]\s*([^\n,\.]+)',
                r'creator\s*[:\-]\s*([^\n,\.]+)',
                # "with Influencer X" or "with Creator X" — name follows the keyword
                r'with\s+((?:influencer|creator)\s+\S+(?:\s+\S+)?)',
                # "with @handle"
                r'with\s+(@\w+)',
            ], user_prompt)
            if val and len(val.strip()) > 1:
                extracted_data["creator"] = val.strip()[:40]

        if not extracted_data.get("startDate"):
            val = _first([r'(?:ran|from|start(?:ed)?)[:\s]+([A-Za-z]+\s+\d+)'], user_prompt)
            if val:
                extracted_data["startDate"] = val

        if not extracted_data.get("endDate"):
            val = _first([r'(?:to|-|through|until|end(?:ed)?)[:\s]+([A-Za-z]+\s+\d+)'], user_prompt)
            if val:
                extracted_data["endDate"] = val

        # Ensure creators array is populated
        if not extracted_data.get("creators") or len(extracted_data["creators"]) == 0:
            print("[NOTE] Populating creators array from extracted data...")
            creator_name = extracted_data.get("creator", "Creator")
            ig = extracted_data.get("instagram", {})
            yt = extracted_data.get("youtube", {})

            extracted_data["creators"] = [
                {
                    "name": creator_name,
                    "platform": "Instagram" if ig.get("views") else "YouTube",
                    "views": ig.get("views", ""),
                    "likes": ig.get("likes", ""),
                    "comments": ig.get("comments", ""),
                    "shares": ig.get("shares", ""),
                    "saves": ig.get("saves", ""),
                    "reach": ig.get("reach", ""),
                    "engagementRate": ig.get("engagementRate", ""),
                    "watchTime": yt.get("watchTime", ""),
                    "interactions": extracted_data.get("performance", {}).get("totalInteractions", ""),
                }
            ]

        print("[OK] Successfully extracted data from images")
        print(f"[DEBUG] Parsed data: campaignName={extracted_data.get('campaignName')}, brand={extracted_data.get('brand')}, ig_views={extracted_data.get('instagram',{}).get('views')}, creators={len(extracted_data.get('creators',[]))}")
        # Ensure extracted data is properly validated
        return validate_and_clean_data(extracted_data)

    except json.JSONDecodeError as e:
        print(f"[ERROR] JSON Parse Error: {str(e)}")
        print(f"Attempted to parse: {json_str if 'json_str' in locals() else response_text}")
        return create_default_data()

    except Exception as e:
        print(f"[ERROR] Error analyzing images with {api_name}: {str(e)}")
        import traceback
        traceback.print_exc()
        return create_default_data()


def populate_powerpoint(data, images):
    """Populate PowerPoint template with extracted data and images"""
    print("[FILE] Creating PowerPoint presentation...")
    
    try:
        # Load template
        prs = Presentation(TEMPLATE_PATH)
        print(f"[OK] Template loaded: {len(prs.slides)} slides")

        def _download_image_stream(url):
            if not url:
                return None
            try:
                import requests
                headers = {
                    "User-Agent": (
                        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                        "AppleWebKit/537.36 (KHTML, like Gecko) "
                        "Chrome/123.0.0.0 Safari/537.36"
                    ),
                    "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
                    "Referer": "https://www.instagram.com/",
                }
                resp = requests.get(url, headers=headers, timeout=30)
                resp.raise_for_status()
                content = resp.content
                content_type = (resp.headers.get("Content-Type") or "").lower()

                # PowerPoint is happiest with PNG/JPEG; convert unsupported formats like WEBP.
                try:
                    img = Image.open(io.BytesIO(content))
                    out = io.BytesIO()
                    save_format = "PNG"
                    if img.mode in ("RGBA", "LA", "P"):
                        img = img.convert("RGBA")
                    else:
                        img = img.convert("RGB")
                        save_format = "JPEG"
                    img.save(out, format=save_format)
                    out.seek(0)
                    print(f"[OK] Downloaded fallback image for PPT ({content_type or 'unknown'})")
                    return out
                except Exception:
                    raw = io.BytesIO(content)
                    raw.seek(0)
                    print(f"[OK] Downloaded raw fallback image for PPT ({content_type or 'unknown'})")
                    return raw
            except Exception as e:
                print(f"[WARNING] Could not download image from URL: {e}")
                return None

        def _uploaded_image_stream(image_obj):
            try:
                img_bytes = image_obj.read()
                image_obj.seek(0)
                return io.BytesIO(img_bytes)
            except Exception as e:
                print(f"[WARNING] Could not read uploaded image: {e}")
                return None

        def _fit_image_within_box(img_stream, box):
            try:
                img_stream.seek(0)
                with Image.open(img_stream) as img:
                    img_width, img_height = img.size
                box_width = box['width']
                box_height = box['height']
                if not img_width or not img_height or not box_width or not box_height:
                    return box['left'], box['top'], box_width, box_height

                image_ratio = img_width / img_height
                box_ratio = box_width / box_height

                if image_ratio > box_ratio:
                    target_width = box_width
                    target_height = int(box_width / image_ratio)
                    target_left = box['left']
                    target_top = box['top'] + int((box_height - target_height) / 2)
                else:
                    target_height = box_height
                    target_width = int(box_height * image_ratio)
                    target_top = box['top']
                    target_left = box['left'] + int((box_width - target_width) / 2)

                return target_left, target_top, target_width, target_height
            except Exception as e:
                print(f"[WARNING] Could not calculate image fit: {e}")
                return box['left'], box['top'], box['width'], box['height']

        def _add_picture_with_fallback(slide, img_stream, placeholder=None, default_box=None, rotation=0):
            if not img_stream:
                return False
            try:
                img_stream.seek(0)
                if placeholder:
                    target_left, target_top, target_width, target_height = _fit_image_within_box(img_stream, placeholder)
                    img_stream.seek(0)
                    picture = slide.shapes.add_picture(
                        img_stream,
                        target_left,
                        target_top,
                        width=target_width,
                        height=target_height
                    )
                    picture.rotation = rotation or placeholder.get('rotation', 0)
                    return True
                if default_box:
                    target_left, target_top, target_width, target_height = _fit_image_within_box(img_stream, default_box)
                    img_stream.seek(0)
                    picture = slide.shapes.add_picture(
                        img_stream,
                        target_left,
                        target_top,
                        width=target_width,
                        height=target_height
                    )
                    picture.rotation = default_box.get('rotation', 0)
                    return True
            except Exception as e:
                print(f"[WARNING] Could not add fallback picture to slide: {e}")
            return False

        def _first_present(*values):
            for value in values:
                if value is None:
                    continue
                if isinstance(value, str):
                    if not value.strip():
                        continue
                    return value
                return value
            return ""

        def _clean_text(value):
            if value is None:
                return ""
            return re.sub(r"\s+", " ", str(value)).strip()

        def _has_value(value):
            cleaned = _clean_text(value)
            return cleaned not in ("", "0", "0.0", "None", "N/A", "NA", "-")

        def _clip_text(value, limit=120):
            cleaned = _clean_text(value)
            if len(cleaned) <= limit:
                return cleaned
            shortened = cleaned[: limit - 1].rsplit(" ", 1)[0].strip()
            return f"{shortened or cleaned[: limit - 1].strip()}…"

        def _display_value(value, fallback="Not available"):
            cleaned = _clean_text(value)
            return cleaned if cleaned else fallback

        def _summarize_sentences(value, sentence_limit=2, char_limit=220):
            cleaned = _clean_text(value)
            if not cleaned:
                return ""
            sentences = re.split(r"(?<=[.!?])\s+", cleaned)
            summary = " ".join(s.strip() for s in sentences if s.strip()) if sentences else cleaned
            if sentences:
                summary = " ".join(s.strip() for s in sentences[:sentence_limit] if s.strip())
            return _clip_text(summary, char_limit)

        def _prepare_text_frame(text_frame):
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        def _pick_title_font_size(text, large=53, medium=44, compact=36, minimum=28):
            length = len(_clean_text(text))
            if length <= 24:
                return Pt(large)
            if length <= 42:
                return Pt(medium)
            if length <= 68:
                return Pt(compact)
            return Pt(minimum)

        def _style_paragraph(paragraph, alignment=None, line_spacing=1.0, space_after=Pt(2)):
            if alignment is not None:
                paragraph.alignment = alignment
            paragraph.line_spacing = line_spacing
            paragraph.space_after = space_after
            paragraph.space_before = Pt(0)

        def _set_title_text(shape, text, *, align=PP_ALIGN.CENTER, italic=False):
            if not hasattr(shape, 'text_frame'):
                shape.text = _clip_text(text, 90)
                return
            text_frame = shape.text_frame
            _prepare_text_frame(text_frame)
            p = text_frame.paragraphs[0]
            p.text = _clip_text(text, 90)
            _style_paragraph(p, alignment=align, line_spacing=0.95, space_after=Pt(0))
            for run in p.runs:
                run.font.name = 'YouTube Sans'
                run.font.size = _pick_title_font_size(text)
                run.font.bold = True
                run.font.italic = italic
                run.font.color.rgb = RGBColor(0x2B, 0x3E, 0x5C)

        def _replace_with_title_box(slide, target_shape, text, *, align=PP_ALIGN.LEFT, italic=True):
            left = getattr(target_shape, 'left', 700000)
            top = getattr(target_shape, 'top', 650000)
            width = getattr(target_shape, 'width', 3000000)
            height = min(getattr(target_shape, 'height', 900000), 950000)

            try:
                sp = target_shape.element
                sp.getparent().remove(sp)
            except Exception:
                pass

            title_box = slide.shapes.add_textbox(left, top, width, height)
            _set_title_text(title_box, text, align=align, italic=italic)
            return title_box

        def _set_textbox_text(slide, left, top, width, height, text, *, font_name='YouTube Sans', font_size=30, bold=True, italic=False, color=(0x2B, 0x3E, 0x5C), align=PP_ALIGN.LEFT):
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            _prepare_text_frame(text_frame)
            paragraph = text_frame.paragraphs[0]
            paragraph.text = _clip_text(text, 120)
            _style_paragraph(paragraph, alignment=align, line_spacing=0.95, space_after=Pt(0))
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.italic = italic
                run.font.color.rgb = RGBColor(*color)
            return text_box

        def _remove_shape(shape):
            try:
                sp = shape.element
                sp.getparent().remove(sp)
            except Exception:
                pass

        def _get_brand_campaign_title():
            campaign_name = _clean_text(data.get('campaignName') or data.get('campaign_name') or '')
            brand_name = _clean_text(data.get('brand') or data.get('brand_name') or '')
            combined_context = f"{campaign_name} {_clean_text(data.get('deliverables'))} {_clean_text(data.get('creator'))}".lower()
            occasion = ""
            for keyword, label in [
                ("diwali", "Diwali"),
                ("summer", "Summer"),
                ("winter", "Winter"),
                ("launch", "Launch"),
                ("airport", "Travel"),
                ("travel", "Travel"),
                ("fashion", "Fashion"),
            ]:
                if keyword in combined_context:
                    occasion = label
                    break
            if campaign_name and brand_name:
                if occasion:
                    return f"{brand_name} {occasion} Campaign"
                if brand_name.lower() in campaign_name.lower() and len(campaign_name) <= 38:
                    return campaign_name
                return f"{brand_name} Creator Campaign"
            if brand_name:
                if occasion:
                    return f"{brand_name} {occasion} Campaign"
                return f"{brand_name} Creator Campaign"
            return campaign_name or 'Campaign Report'

        def _set_metric_block(text_frame, metrics, extras=None, learnings_text=""):
            _prepare_text_frame(text_frame)
            visible_metrics = [(label, _clean_text(value)) for label, value in (metrics or []) if _has_value(value)]
            visible_extras = [(label, _clip_text(value, 72)) for label, value in (extras or []) if _has_value(value)]
            learnings_summary = _summarize_sentences(learnings_text, sentence_limit=2, char_limit=190)

            total_lines = len(visible_metrics) + len(visible_extras) + (1 if learnings_summary else 0)
            body_size = 16
            if total_lines >= 14:
                body_size = 11
            elif total_lines >= 12:
                body_size = 12
            elif total_lines >= 10:
                body_size = 13
            elif total_lines >= 8:
                body_size = 14
            elif total_lines >= 6:
                body_size = 15

            def _add_label_value(paragraph, label, value):
                _style_paragraph(paragraph, line_spacing=1.0 if body_size >= 15 else 0.96, space_after=Pt(1))
                label_run = paragraph.add_run()
                label_run.text = label
                label_run.font.bold = True
                label_run.font.name = 'Aptos'
                label_run.font.size = Pt(body_size)
                label_run.font.color.rgb = RGBColor(0x2B, 0x3E, 0x5C)

                value_run = paragraph.add_run()
                value_run.text = value
                value_run.font.bold = False
                value_run.font.name = 'Aptos'
                value_run.font.size = Pt(body_size)
                value_run.font.color.rgb = RGBColor(0x2B, 0x3E, 0x5C)

            for index, (label, value) in enumerate(visible_metrics):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                _add_label_value(paragraph, label, value)

            line_capacity = 10
            if total_lines > line_capacity:
                line_capacity = min(total_lines, 15)
            remaining_slots = max(0, line_capacity - len(visible_metrics) - (1 if learnings_summary else 0))
            for label, value in visible_extras[:remaining_slots]:
                paragraph = text_frame.add_paragraph()
                _add_label_value(paragraph, label, value)

            if learnings_summary:
                paragraph = text_frame.add_paragraph()
                _style_paragraph(paragraph, line_spacing=1.0 if body_size >= 15 else 0.96, space_after=Pt(0))

                label_run = paragraph.add_run()
                label_run.text = "Learnings: "
                label_run.font.bold = True
                label_run.font.name = 'Aptos'
                label_run.font.size = Pt(max(body_size - 1, 12))
                label_run.font.color.rgb = RGBColor(0x2B, 0x3E, 0x5C)

                value_run = paragraph.add_run()
                value_run.text = learnings_summary
                value_run.font.bold = False
                value_run.font.name = 'Aptos'
                value_run.font.size = Pt(max(body_size - 1, 12))
                value_run.font.color.rgb = RGBColor(0x2B, 0x3E, 0x5C)

        def _insert_new_slide_at(presentation, layout, insert_index):
            new_slide = presentation.slides.add_slide(layout)
            slide_id_list = presentation.slides._sldIdLst
            new_slide_id = slide_id_list[-1]
            slide_id_list.remove(new_slide_id)
            slide_id_list.insert(insert_index, new_slide_id)
            return new_slide

        def _render_campaign_gallery_slide(slide, image_items, page_index, total_pages):
            title_box = slide.shapes.add_textbox(260000, 140000, 8600000, 520000)
            title_tf = title_box.text_frame
            _prepare_text_frame(title_tf)
            title_para = title_tf.paragraphs[0]
            _style_paragraph(title_para, alignment=PP_ALIGN.LEFT, line_spacing=1.0, space_after=Pt(0))
            title_run = title_para.add_run()
            title_run.text = f"Campaign Media Gallery ({page_index}/{total_pages})"
            title_run.font.name = 'Aptos'
            title_run.font.bold = True
            title_run.font.size = Pt(24)
            title_run.font.color.rgb = RGBColor(0x2B, 0x3E, 0x5C)

            grid_boxes = [
                {'left': 260000,  'top': 760000,  'width': 2750000, 'height': 1650000},
                {'left': 3190000, 'top': 760000,  'width': 2750000, 'height': 1650000},
                {'left': 6120000, 'top': 760000,  'width': 2750000, 'height': 1650000},
                {'left': 260000,  'top': 2520000, 'width': 2750000, 'height': 1650000},
                {'left': 3190000, 'top': 2520000, 'width': 2750000, 'height': 1650000},
                {'left': 6120000, 'top': 2520000, 'width': 2750000, 'height': 1650000},
            ]

            for idx, image_item in enumerate(image_items[:6]):
                if idx >= len(grid_boxes):
                    break
                stream_obj = image_item.get("stream")
                if not stream_obj:
                    continue
                box = grid_boxes[idx]
                _add_picture_with_fallback(
                    slide,
                    stream_obj,
                    default_box={**box, 'rotation': 0},
                )

                caption_text = image_item.get("label") or f"Image {idx + 1}"
                caption_box = slide.shapes.add_textbox(
                    box['left'],
                    box['top'] + box['height'] + 60000,
                    box['width'],
                    280000,
                )
                caption_tf = caption_box.text_frame
                _prepare_text_frame(caption_tf)
                caption_para = caption_tf.paragraphs[0]
                _style_paragraph(caption_para, alignment=PP_ALIGN.LEFT, line_spacing=1.0, space_after=Pt(0))
                caption_run = caption_para.add_run()
                caption_run.text = _clip_text(caption_text, 62)
                caption_run.font.name = 'Aptos'
                caption_run.font.bold = False
                caption_run.font.size = Pt(10)
                caption_run.font.color.rgb = RGBColor(0x2B, 0x3E, 0x5C)

        def _dedupe_creators_for_slides(creators):
            deduped = []
            key_index = {}

            for item in creators:
                if not isinstance(item, dict):
                    continue

                normalized = dict(item)
                name_value = (normalized.get("name") or normalized.get("creator_name") or "").strip()
                post_url_value = (normalized.get("postUrl") or "").strip().lower()
                name_key = re.sub(r"[^a-z0-9]+", "", name_value.lower())
                dedupe_key = post_url_value or name_key or f"creator_{len(deduped)}"

                if dedupe_key in key_index:
                    existing_item = deduped[key_index[dedupe_key]]
                    for merge_key, merge_value in normalized.items():
                        if merge_value in (None, "", "N/A"):
                            continue
                        if existing_item.get(merge_key) in (None, "", "N/A"):
                            existing_item[merge_key] = merge_value
                    continue

                key_index[dedupe_key] = len(deduped)
                deduped.append(normalized)

            return deduped

        def _clone_slide_from_template(presentation, source_slide, insert_index=None):
            new_slide = presentation.slides.add_slide(source_slide.slide_layout)

            for shape in list(new_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)

            for shape_element in source_slide.shapes._spTree:
                if shape_element.tag.endswith('extLst'):
                    continue
                if shape_element.tag.endswith('pic'):
                    continue
                new_slide.shapes._spTree.insert_element_before(copy.deepcopy(shape_element), 'p:extLst')

            for shape in source_slide.shapes:
                if getattr(shape, 'shape_type', None) == 13:
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        picture = new_slide.shapes.add_picture(
                            image_stream,
                            shape.left,
                            shape.top,
                            width=shape.width,
                            height=shape.height
                        )
                        picture.rotation = getattr(shape, 'rotation', 0)
                    except Exception as picture_copy_error:
                        print(f"[WARNING] Could not clone picture shape: {picture_copy_error}")

            if insert_index is not None:
                slide_id_list = presentation.slides._sldIdLst
                new_slide_id = slide_id_list[-1]
                slide_id_list.remove(new_slide_id)
                slide_id_list.insert(insert_index, new_slide_id)

            return new_slide

        def _render_creator_detail_slide(slide, creator_item, creator_index, creator_total, creator_image_stream=None, insight_stream=None):
            creator_name_local = creator_item.get('name', 'Creator')
            creator_budget_display_local = creator_item.get('budget') or data.get('financial', {}).get('totalBudget') or ''
            budget_currency_local = data.get('financial', {}).get('budgetCurrency') or ''
            if creator_budget_display_local and budget_currency_local:
                creator_budget_display_local = f"{creator_budget_display_local} {budget_currency_local}"

            brand_name_local = creator_item.get('brand') or data.get('brand') or data.get('brand_name', 'Brand')
            for shape in list(slide.shapes):
                shape_text = getattr(shape, 'text', '')
                if not hasattr(shape, 'text_frame'):
                    continue
                if 'Confidential' in shape_text:
                    continue
                if getattr(shape, 'left', 0) < 3600000 and getattr(shape, 'top', 0) < 4300000:
                    _remove_shape(shape)

            _set_textbox_text(
                slide,
                430000,
                430000,
                3000000,
                430000,
                f"{creator_name_local} for {brand_name_local}",
                font_size=18,
                bold=True,
                italic=True,
                align=PP_ALIGN.LEFT,
            )

            financial_data_local = data.get('financial', {})
            metrics_local = [
                ("No. of Views: ", _display_value(_first_present(creator_item.get('views')))),
                ("Likes: ", _display_value(_first_present(creator_item.get('likes')))),
                ("Comments: ", _display_value(_first_present(creator_item.get('comments')))),
                ("Shares: ", _display_value(_first_present(creator_item.get('shares')))),
                ("Saves: ", _display_value(_first_present(creator_item.get('saves')))),
                ("Reach: ", _display_value(_first_present(creator_item.get('reach')))),
                ("Engagement Rate: ", _display_value(_first_present(creator_item.get('engagementRate')))),
                ("Total Engagement: ", _display_value(_prefer_interactions(creator_item.get('total_engagement'), creator_item.get('interactions')))),
                ("Budget: ", _display_value(creator_budget_display_local)),
                ("CPC: ", _display_value(_first_present(creator_item.get('cpc'), financial_data_local.get('cpc')))),
                ("CPV: ", _display_value(_first_present(creator_item.get('cpv'), financial_data_local.get('cpv')))),
                ("CPE: ", _display_value(_first_present(creator_item.get('cpe'), financial_data_local.get('cpe'))))
            ]
            extras_local = [
                ("Slide: ", f"{creator_index + 1} of {creator_total}"),
                ("Platform: ", _display_value(creator_item.get('platform') or 'Instagram')),
                # ("CPC Goal: ", _display_value(creator_item.get('cpcGoal') or financial_data_local.get('cpcGoal') or '')),
                # ("CPV Goal: ", _display_value(creator_item.get('cpvGoal') or financial_data_local.get('cpvGoal') or '')),
                # ("CPC Calculation: ", _display_value(creator_item.get('cpcCalculation') or financial_data_local.get('cpcCalculation') or '')),
                # ("CPV Calculation: ", _display_value(creator_item.get('cpvCalculation') or financial_data_local.get('cpvCalculation') or '')),
            ]
            learnings_text_local = creator_item.get('learnings', data.get('performance', {}).get('keyLearnings', ''))
            metric_box = slide.shapes.add_textbox(260000, 1470000, 2480000, 1980000)
            _set_metric_block(metric_box.text_frame, metrics_local, extras_local, learnings_text_local)

            if not creator_image_stream:
                creator_image_stream = _download_image_stream(creator_item.get("postImage") or data.get("postImage"))

            if creator_image_stream:
                text_placeholder = None
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and 'Creator content picture here' in shape.text:
                        text_placeholder = shape
                        break
                if text_placeholder:
                    sp = text_placeholder.element
                    sp.getparent().remove(sp)

                creator_placeholder = None
                for shape in slide.shapes:
                    if shape.shape_type == 1 and 3900000 < shape.left < 4500000 and 1000000 < shape.top < 2000000:
                        creator_placeholder = {
                            'shape': shape,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        }
                        break

                if creator_placeholder:
                    sp = creator_placeholder['shape'].element
                    sp.getparent().remove(sp)
                    _add_picture_with_fallback(slide, creator_image_stream, placeholder=creator_placeholder)
                else:
                    default_creator_box = {
                        'left': 3958173,
                        'top': 1105651,
                        'width': 2559000,
                        'height': 2932200,
                        'rotation': 0,
                    }
                    _add_picture_with_fallback(slide, creator_image_stream, default_box=default_creator_box)

            if insight_stream:
                text_placeholder = None
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and 'content insights here' in shape.text:
                        text_placeholder = shape
                        break
                if text_placeholder:
                    sp = text_placeholder.element
                    sp.getparent().remove(sp)

                insight_placeholder = None
                for shape in slide.shapes:
                    if shape.shape_type == 1 and shape.left > 6000000:
                        insight_placeholder = {
                            'shape': shape,
                            'left': shape.left,
                            'top': shape.top,
                            'width': shape.width,
                            'height': shape.height
                        }
                        break

                if insight_placeholder:
                    sp = insight_placeholder['shape'].element
                    sp.getparent().remove(sp)
                    insight_stream.seek(0)
                    slide.shapes.add_picture(
                        insight_stream,
                        insight_placeholder['left'],
                        insight_placeholder['top'],
                        width=insight_placeholder['width'],
                        height=insight_placeholder['height']
                    )

        def _prefer_interactions(total_engagement_value, total_interactions_value):
            interaction_value = _first_present(total_interactions_value)
            engagement_value = _first_present(total_engagement_value)
            if interaction_value not in ("", None):
                return interaction_value
            return engagement_value

        image_classifications = data.get('image_classifications', []) or []
        type_to_indices = {}
        for classification in image_classifications:
            idx = classification.get('image_index')
            img_type = classification.get('type')
            if isinstance(idx, int) and img_type:
                type_to_indices.setdefault(img_type, []).append(idx)

        used_uploaded_indices = set()

        def _select_uploaded_streams(preferred_types=None, limit=1, fallback_any=False):
            selected = []
            selected_indices = []
            preferred_types = preferred_types or []

            for img_type in preferred_types:
                for idx in type_to_indices.get(img_type, []):
                    if idx in used_uploaded_indices or idx >= len(images):
                        continue
                    stream = _uploaded_image_stream(images[idx])
                    if stream:
                        selected.append(stream)
                        selected_indices.append(idx)
                        used_uploaded_indices.add(idx)
                        if len(selected) >= limit:
                            return selected, selected_indices

            if fallback_any and len(selected) < limit:
                for idx, image_obj in enumerate(images):
                    if idx in used_uploaded_indices:
                        continue
                    stream = _uploaded_image_stream(image_obj)
                    if stream:
                        selected.append(stream)
                        selected_indices.append(idx)
                        used_uploaded_indices.add(idx)
                        if len(selected) >= limit:
                            break

            return selected, selected_indices

        def _find_logo_placeholder(slide):
            fallback_logo_placeholder = None
            for shape in slide.shapes:
                if not hasattr(shape, 'text'):
                    continue
                shape_text = (shape.text or "").strip().upper()
                if not shape_text:
                    continue
                if "BRAND LOGO HERE" in shape_text:
                    return shape
                if "LOGO" in shape_text and fallback_logo_placeholder is None:
                    fallback_logo_placeholder = shape
            return fallback_logo_placeholder

        def _resolve_brand_logo_bytes(logo_images):
            if logo_images:
                img_bytes = logo_images[0].read()
                logo_images[0].seek(0)
                if img_bytes:
                    return img_bytes, "uploaded_image"

            brand_logo_url = (data.get('brandLogo') or '').strip()
            if brand_logo_url:
                brand_logo_stream = _download_image_stream(brand_logo_url)
                if brand_logo_stream:
                    return brand_logo_stream.getvalue(), "brand_logo_url"
            return b"", ""

        def _place_brand_logo(slide, logo_bytes, slide_label, fallback_box):
            if not logo_bytes:
                return False

            logo_placeholder = _find_logo_placeholder(slide)
            box = dict(fallback_box)
            if logo_placeholder:
                box = {
                    'left': logo_placeholder.left,
                    'top': logo_placeholder.top,
                    'width': logo_placeholder.width,
                    'height': logo_placeholder.height,
                    'rotation': 0,
                }
                _remove_shape(logo_placeholder)

            with io.BytesIO(logo_bytes) as logo_stream:
                _add_picture_with_fallback(
                    slide,
                    logo_stream,
                    default_box=box,
                )
            print(f"[OK] Brand logo inserted on {slide_label}")
            return True
        
        # SLIDE 1: Insert Brand Logo
        hero_image_stream = None
        brand_logo_bytes = b""
        if len(prs.slides) > 0:
            slide_1 = prs.slides[0]
            
            # Find and replace "BRAND LOGO HERE" text box with brand logo image
            logo_images = [
                img for idx, img in enumerate(images) 
                if any(c.get('image_index') == idx and c.get('type') == 'brand_logo' 
                      for c in data.get('image_classifications', []))
            ]
            
            brand_logo_bytes, logo_source = _resolve_brand_logo_bytes(logo_images)
            if brand_logo_bytes:
                print(f"[DESIGN] Adding brand logo from {logo_source}")
                _place_brand_logo(
                    slide_1,
                    brand_logo_bytes,
                    "Slide 1",
                    fallback_box={
                        'left': 4675100,
                        'top': 2084875,
                        'width': 2552100,
                        'height': 461700,
                        'rotation': 0,
                    }
                )
            else:
                print("[WARNING] No brand logo source found for Slide 1")

            slide_1_image_streams, _ = _select_uploaded_streams(
                preferred_types=['campaign_photo', 'creator_content'],
                limit=1,
                fallback_any=False,
            )
            if not slide_1_image_streams:
                slide_1_image_url = ""
                for creator_item in data.get("creators", []) or []:
                    if not isinstance(creator_item, dict):
                        continue
                    slide_1_image_url = (creator_item.get("postImage") or "").strip()
                    if slide_1_image_url:
                        break
                if not slide_1_image_url:
                    slide_1_image_url = (data.get("postImage") or "").strip()
                if slide_1_image_url:
                    fallback_slide_1_image = _download_image_stream(slide_1_image_url)
                    if fallback_slide_1_image:
                        slide_1_image_streams.append(fallback_slide_1_image)

            if slide_1_image_streams:
                hero_image_stream = slide_1_image_streams[0]
                for shape in list(slide_1.shapes):
                    if shape.shape_type == 1 and getattr(shape, 'top', 0) > 3200000:
                        _remove_shape(shape)
                _add_picture_with_fallback(
                    slide_1,
                    hero_image_stream,
                    default_box={
                        'left': 1700000,
                        'top': 3200000,
                        'width': 5250000,
                        'height': 1180000,
                        'rotation': 0,
                    }
                )
                print("[OK] Added dynamic campaign image to Slide 1")

        # SLIDE 2: Add dynamic campaign image similar to the opening hero treatment
        if len(prs.slides) > 1 and hero_image_stream:
            slide_2 = prs.slides[1]
            if brand_logo_bytes:
                _place_brand_logo(
                    slide_2,
                    brand_logo_bytes,
                    "Slide 2",
                    fallback_box={
                        'left': 6620000,
                        'top': 220000,
                        'width': 2200000,
                        'height': 520000,
                        'rotation': 0,
                    }
                )
            for shape in list(slide_2.shapes):
                if shape.shape_type == 1 and getattr(shape, 'left', 0) < 3600000 and getattr(shape, 'top', 0) > 900000:
                    _remove_shape(shape)
            _add_picture_with_fallback(
                slide_2,
                hero_image_stream,
                default_box={
                    'left': 250000,
                    'top': 1250000,
                    'width': 3600000,
                    'height': 2900000,
                    'rotation': 0,
                }
            )
            print("[OK] Added dynamic campaign image to Slide 2")
        
        # SLIDE 3: Campaign Name with proper styling
        if len(prs.slides) > 2:
            slide_3 = prs.slides[2]
            
            for shape in slide_3.shapes:
                if hasattr(shape, 'text') and '[Campaign Name Here]' in shape.text:
                    if hasattr(shape, 'text_frame'):
                        campaign_name = _get_brand_campaign_title()
                        _set_title_text(shape, campaign_name, align=PP_ALIGN.CENTER, italic=False)

                        print(f"[OK] Updated Slide 3 with styled campaign name: {campaign_name}")
        
        campaign_gallery_items = []

        # SLIDE 4: Overall Campaign Report with bold data
        if len(prs.slides) > 3:
            slide_4 = prs.slides[3]
            # Support both old and new data formats
            overall = data.get('overall_campaign', {})
            instagram = data.get('instagram', {})
            financial = data.get('financial', {})
            performance = data.get('performance', {})
            creator_items = [item for item in data.get("creators", []) if isinstance(item, dict)]
            unique_creator_names = []
            seen_creator_names = set()
            for creator_item in creator_items:
                creator_name_candidate = (creator_item.get("name") or "").strip()
                if not creator_name_candidate:
                    continue
                creator_name_key = creator_name_candidate.lower()
                if creator_name_key in seen_creator_names:
                    continue
                seen_creator_names.add(creator_name_key)
                unique_creator_names.append(creator_name_candidate)

            creators_display = ", ".join(unique_creator_names[:4]) if unique_creator_names else _display_value(data.get('creator') or '')
            budget_display = financial.get('totalBudget') or overall.get('budget') or ''
            if budget_display and financial.get('budgetCurrency'):
                budget_display = f"{budget_display} {financial.get('budgetCurrency')}"

            for shape in slide_4.shapes:
                if hasattr(shape, 'text_frame') and hasattr(shape, 'text') and 'No. of Views:' in shape.text:
                    text_frame = shape.text_frame

                    # Define metrics with labels and values (support both old and new formats)
                    metrics = [
                        ("No. of Views: ", _display_value(_first_present(performance.get('totalViews'), instagram.get('views'), overall.get('views')))),
                        ("Likes: ", _display_value(_first_present(performance.get('totalLikes'), instagram.get('likes'), overall.get('likes')))),
                        ("Comments: ", _display_value(_first_present(performance.get('totalComments'), instagram.get('comments'), overall.get('comments')))),
                        ("Shares: ", _display_value(_first_present(performance.get('totalShares'), instagram.get('shares'), overall.get('shares')))),
                        ("Saves: ", _display_value(_first_present(performance.get('totalSaves'), instagram.get('saves'), overall.get('saves')))),
                        ("Reach: ", _display_value(_first_present(performance.get('totalReach'), instagram.get('reach'), overall.get('reach')))),
                        ("Engagement Rate: ", _display_value(_first_present(instagram.get('engagementRate'), overall.get('engagementRate')))),
                        ("Total Engagement: ", _display_value(_prefer_interactions(performance.get('totalEngagement'), performance.get('totalInteractions')) or overall.get('total_engagement') or '')),
                        ("Budget: ", _display_value(budget_display)),
                    ]

                    extras = [
                        ("Campaign: ", _display_value(data.get('campaignName') or data.get('campaign_name') or '')),
                        ("Brand: ", _display_value(data.get('brand') or data.get('brand_name') or '')),
                        ("Creators: ", _display_value(creators_display)),
                        ("Creator Count: ", _display_value(str(len(unique_creator_names)) if unique_creator_names else '')),
                        ("CPC: ", _display_value(_first_present(financial.get('cpc'), overall.get('cpc')))),
                        ("CPV: ", _display_value(_first_present(financial.get('cpv'), overall.get('cpv')))),
                        ("CPE: ", _display_value(_first_present(financial.get('cpe'), overall.get('cpe')))),
                        # ("CPC Goal: ", _display_value(financial.get('cpcGoal') or '')),
                        # ("CPV Goal: ", _display_value(financial.get('cpvGoal') or '')),
                        # ("CPC Calculation: ", _display_value(financial.get('cpcCalculation') or '')),
                        # ("CPV Calculation: ", _display_value(financial.get('cpvCalculation') or '')),
                    ]
                    learnings_text = (
                        overall.get('learnings')
                        or performance.get('keyLearnings')
                        or 'No learnings available.'
                    )
                    _set_metric_block(text_frame, metrics, extras, learnings_text)
                    
                    print("[OK] Updated Slide 4 with bold formatting for overall campaign metrics")
                    break
            
            # Add campaign photos (up to 3) - Replace template placeholders with rotation
            campaign_image_streams, campaign_image_indices = _select_uploaded_streams(
                preferred_types=['campaign_photo'],
                limit=10,
                fallback_any=False,
            )
            campaign_image_items = []
            for idx, stream in enumerate(campaign_image_streams):
                upload_idx = campaign_image_indices[idx] if idx < len(campaign_image_indices) else None
                label = f"Uploaded image {idx + 1}" if upload_idx is None else f"Uploaded image #{upload_idx + 1}"
                campaign_image_items.append({"stream": stream, "label": label})
            if not campaign_image_streams:
                url_image_candidates = []
                seen_url_images = set()

                for creator_item in data.get("creators", []) or []:
                    if not isinstance(creator_item, dict):
                        continue
                    image_url = (creator_item.get("postImage") or "").strip()
                    if image_url and image_url not in seen_url_images:
                        seen_url_images.add(image_url)
                        url_image_candidates.append(image_url)
                    if len(url_image_candidates) >= 10:
                        break

                top_level_post_image = (data.get("postImage") or "").strip()
                if top_level_post_image and top_level_post_image not in seen_url_images and len(url_image_candidates) < 10:
                    seen_url_images.add(top_level_post_image)
                    url_image_candidates.append(top_level_post_image)

                for image_url in url_image_candidates[:10]:
                    fallback_campaign_image = _download_image_stream(image_url)
                    if fallback_campaign_image:
                        campaign_image_streams.append(fallback_campaign_image)
                        campaign_image_items.append({"stream": fallback_campaign_image, "label": image_url})
            elif not campaign_image_items:
                for idx, stream in enumerate(campaign_image_streams):
                    campaign_image_items.append({"stream": stream, "label": f"Campaign image {idx + 1}"})
            
            if campaign_image_streams:
                print(f"[IMAGE] Adding {len(campaign_image_streams[:3])} campaign photos to Slide 4")
                
                # Find and delete text placeholder "Campaign Photos Here at Least 3 photos"
                text_placeholder = None
                shapes_to_check = list(slide_4.shapes)  # Create a list to avoid iteration issues
                for shape in shapes_to_check:
                    if hasattr(shape, 'text') and 'Campaign Photos Here' in shape.text:
                        text_placeholder = shape
                        break
                
                if text_placeholder:
                    sp = text_placeholder.element
                    sp.getparent().remove(sp)
                    print("[OK] Removed 'Campaign Photos Here' text placeholder")
                
                shape_placeholders = []
                shapes_to_check = list(slide_4.shapes)
                for shape in shapes_to_check:
                    if shape.shape_type == 1 and shape.left > 3000000 and shape.top > 900000 and shape.top < 4200000:
                        shape_placeholders.append(shape)

                for placeholder_shape in shape_placeholders:
                    _remove_shape(placeholder_shape)

                layout_boxes_by_count = {
                    1: [
                        {'left': 4800000, 'top': 980000, 'width': 2200000, 'height': 2450000, 'rotation': 0},
                    ],
                    2: [
                        {'left': 4100000, 'top': 980000, 'width': 1650000, 'height': 2300000, 'rotation': -9},
                        {'left': 6350000, 'top': 980000, 'width': 1650000, 'height': 2300000, 'rotation': 5},
                    ],
                    3: [
                        {'left': 3600000, 'top': 1020000, 'width': 1450000, 'height': 2100000, 'rotation': -11},
                        {'left': 5350000, 'top': 900000, 'width': 1650000, 'height': 2350000, 'rotation': 0},
                        {'left': 7350000, 'top': 1020000, 'width': 1350000, 'height': 2050000, 'rotation': 8},
                    ],
                }
                active_layout = layout_boxes_by_count.get(min(len(campaign_image_streams[:3]), 3), layout_boxes_by_count[1])
                for idx, (img_stream, box) in enumerate(zip(campaign_image_streams[:3], active_layout)):
                    _add_picture_with_fallback(slide_4, img_stream, default_box=box, rotation=box.get('rotation', 0))
                    print(f"[OK] Added campaign photo {idx + 1} using fixed layout box")

                if len(campaign_image_items) > 3 and len(prs.slides) > 4:
                    campaign_gallery_items = campaign_image_items[3:10]
                    print(f"[IMAGE] Queued {len(campaign_gallery_items)} additional campaign image(s) for gallery slides")
        
        # SLIDE 5: Creator x Brand with bold data
        if len(prs.slides) > 4:
            slide_5_template = prs.slides[4]

            if campaign_gallery_items:
                chunks = [
                    campaign_gallery_items[idx:idx + 6]
                    for idx in range(0, len(campaign_gallery_items), 6)
                ]
                insert_index = 4
                gallery_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
                for page_idx, chunk in enumerate(chunks, start=1):
                    gallery_slide = _insert_new_slide_at(prs, gallery_layout, insert_index)
                    _render_campaign_gallery_slide(gallery_slide, chunk, page_idx, len(chunks))
                    print(f"[OK] Rendered campaign gallery slide {page_idx}/{len(chunks)} with {len(chunk)} image(s)")
                    insert_index += 1

                # Template moved after gallery insertion; refresh handle at current index.
                slide_5_template = prs.slides[insert_index]

            creators = [item for item in data.get('creators', []) if isinstance(item, dict)]
            if not creators:
                fallback_creator = data.get('creator_data', {}) or {}
                fallback_creator["name"] = fallback_creator.get('creator_name', data.get('creator', 'Creator'))
                creators = [fallback_creator]
            creators = _dedupe_creators_for_slides(creators)

            creator_slides = creators[:10]
            creator_image_streams, _ = _select_uploaded_streams(
                preferred_types=['creator_content'],
                limit=max(1, len(creator_slides)),
                fallback_any=True,
            )
            insight_streams, _ = _select_uploaded_streams(
                preferred_types=['insights_screenshot', 'campaign_dashboard'],
                limit=max(1, len(creator_slides)),
                fallback_any=True,
            )

            creator_detail_slides = [slide_5_template]
            template_slide_id = slide_5_template.slide_id
            template_index = 0
            for idx, slide_obj in enumerate(prs.slides):
                if slide_obj.slide_id == template_slide_id:
                    template_index = idx
                    break
            insert_index = template_index + 1
            for _ in creator_slides[1:]:
                creator_detail_slides.append(
                    _clone_slide_from_template(prs, slide_5_template, insert_index=insert_index)
                )
                insert_index += 1

            for creator_index, (detail_slide, creator_item) in enumerate(zip(creator_detail_slides, creator_slides)):
                _render_creator_detail_slide(
                    detail_slide,
                    creator_item,
                    creator_index,
                    len(creator_slides),
                    creator_image_stream=creator_image_streams[creator_index] if creator_index < len(creator_image_streams) else None,
                    insight_stream=insight_streams[creator_index] if creator_index < len(insight_streams) else (insight_streams[0] if insight_streams else None),
                )
                print(f"[OK] Rendered Slide 5 detail slide {creator_index + 1}/{len(creator_slides)}: {creator_item.get('name', 'Creator')}")
        
        # Save the presentation
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M")
        # Handle both old format (campaign_name) and new format (campaignName)
        campaign_name = data.get('campaignName') or data.get('campaign_name') or 'Campaign Report'
        normalized_campaign_name = re.sub(r'[^A-Za-z0-9]+', '_', str(campaign_name)).strip('_')
        safe_campaign_name = re.sub(r'_+', '_', normalized_campaign_name) or "campaign_report"
        filename = f"PyroMedia_Report_{safe_campaign_name}_{timestamp}.pptx"
        output_path = os.path.join(OUTPUT_DIR, filename)
        
        prs.save(output_path)
        print(f"[OK] Report saved: {output_path}")
        
        return output_path
        
    except Exception as e:
        print(f"[ERROR] Error populating PowerPoint: {str(e)}")
        import traceback
        traceback.print_exc()
        raise


def upload_to_google_drive(file_path, filename):
    """Upload PPTX to Google Drive and convert to Google Slides - PRODUCTION READY (OAuth)"""
    try:
        from google.oauth2.credentials import Credentials
        from google.auth.transport.requests import Request
        from googleapiclient.discovery import build
        from googleapiclient.http import MediaFileUpload
        import pickle
        
        SCOPES = ['https://www.googleapis.com/auth/drive.file']
        
        creds = None
        
        # Try to load from token.pickle file
        if os.path.exists('token.pickle'):
            print("[INFO] Loading credentials from token.pickle...")
            with open('token.pickle', 'rb') as token:
                creds = pickle.load(token)
        else:
            print("[ERROR] ERROR: token.pickle not found!")
            print("Please generate token.pickle by running generate_token.py locally")
            return None
        
        # Refresh token if expired (NO BROWSER NEEDED)
        if creds and creds.expired and creds.refresh_token:
            print("[INFO] Refreshing expired token...")
            creds.refresh(Request())
            
            # Save refreshed token
            with open('token.pickle', 'wb') as token:
                pickle.dump(creds, token)
            print("[OK] Token refreshed successfully")
        
        # Check if credentials are valid
        if not creds or not creds.valid:
            print("[ERROR] ERROR: Invalid credentials!")
            print("Please regenerate token.pickle by running generate_token.py locally")
            return None
        
        # Build Drive API service
        service = build('drive', 'v3', credentials=creds)
        
        # Upload file with conversion to Google Slides
        file_metadata = {
            'name': filename.replace('.pptx', ''),
            'mimeType': 'application/vnd.google-apps.presentation'
        }
        
        media = MediaFileUpload(
            file_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            resumable=True
        )
        
        print("[INFO] Uploading to Google Drive...")
        file = service.files().create(
            body=file_metadata,
            media_body=media,
            fields='id, webViewLink'
        ).execute()
        
        file_link = file.get('webViewLink')
        print(f"[OK] Uploaded to Google Drive: {file_link}")
        
        # Make file accessible to anyone with link
        try:
            permission = {
                'type': 'anyone',
                'role': 'writer',
                'allowFileDiscovery': False
            }
            service.permissions().create(
                fileId=file.get('id'),
                body=permission
            ).execute()
            print(f"[OK] File is now accessible with link")
        except Exception as perm_error:
            print(f"[WARNING]  Could not set link sharing: {perm_error}")
        
        return file_link
        
    except Exception as e:
        print(f"[ERROR] Error uploading to Google Drive: {str(e)}")
        import traceback
        traceback.print_exc()
        return None


@app.route('/')
def home():
    """Home endpoint"""
    return jsonify({
        "message": "Campaign Report Generator API",
        "status": "running",
        "version": "1.0",
        "endpoints": {
            "generate_report": "/api/generate-report (POST)",
            "download": "/api/download/<filename> (GET)",
            "prompt_presets": "/api/prompt-presets (GET)",
            "build_prompt": "/api/build-prompt (POST)",
        }
    })


@app.route('/api/health', methods=['GET'])
def health_check():
    """Health/debug endpoint for deployment verification without exposing secrets."""
    return jsonify({
        "success": True,
        "status": "ok",
        "runtime": {
            "platform": os.name,
            "python_version": sys.version.split()[0],
            "preview_export_supported": os.name == 'nt' and shutil.which("powershell") is not None,
            "using_openrouter": USE_OPENROUTER,
        },
        "env": {
            "openrouter_api_key_present": bool(os.environ.get("OPENROUTER_API_KEY", "").strip()),
            "anthropic_api_key_present": bool(os.environ.get("ANTHROPIC_API_KEY", "").strip()),
            "brightdata_api_token_present": bool(
                os.environ.get("BRIGHTDATA_API_TOKEN", "").strip()
                or os.environ.get("BRIGHTDATA_API_KEY", "").strip()
            ),
            "brightdata_instagram_post_dataset_id_present": bool(
                os.environ.get("BRIGHTDATA_INSTAGRAM_POST_DATASET_ID", "").strip()
            ),
            "brightdata_instagram_reels_dataset_id_present": bool(
                os.environ.get("BRIGHTDATA_INSTAGRAM_REELS_DATASET_ID", "").strip()
            ),
            "brightdata_instagram_posts_dataset_id_present": bool(
                os.environ.get("BRIGHTDATA_INSTAGRAM_POSTS_DATASET_ID", "").strip()
            ),
            "brightdata_instagram_profile_dataset_id_present": bool(
                os.environ.get("BRIGHTDATA_INSTAGRAM_PROFILE_DATASET_ID", "").strip()
                or os.environ.get("BRIGHTDATA_INSTAGRAM_DATASET_ID", "").strip()
            ),
        },
    })


@app.route('/api/openapi.json', methods=['GET'])
def openapi_spec():
    """Minimal OpenAPI spec for interactive backend docs."""
    server_url = request.host_url.rstrip('/')
    spec = {
        "openapi": "3.0.3",
        "info": {
            "title": "Phyo Report Generator API",
            "version": "1.0.0",
            "description": "Interactive docs for the campaign report generator backend.",
        },
        "servers": [{"url": server_url}],
        "paths": {
            "/api/health": {
                "get": {
                    "summary": "Health check",
                    "responses": {
                        "200": {
                            "description": "Deployment/runtime health information",
                        }
                    },
                }
            },
            "/api/prompt-presets": {
                "get": {
                    "summary": "Get prompt presets",
                    "responses": {"200": {"description": "Prompt templates"}},
                }
            },
            "/api/build-prompt": {
                "post": {
                    "summary": "Build prompt from template",
                    "requestBody": {
                        "required": True,
                        "content": {
                            "application/json": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "template": {"type": "string", "example": "master"},
                                        "data": {"type": "object"},
                                    },
                                }
                            }
                        },
                    },
                    "responses": {
                        "200": {"description": "Prompt built successfully"},
                        "400": {"description": "Invalid template"},
                    },
                }
            },
            "/api/generate-report": {
                "post": {
                    "summary": "Generate campaign report",
                    "requestBody": {
                        "required": True,
                        "content": {
                            "multipart/form-data": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "prompt": {"type": "string"},
                                        "instagram_post_url": {"type": "string", "description": "One or more URLs separated by new lines (up to 10 processed)"},
                                        "budget_inr_values": {
                                            "type": "string",
                                            "description": "JSON string array for URL-wise budgets (up to 10 entries)",
                                            "example": "[\"10000\",\"20000\",\"\",\"\",\"\"]"
                                        },
                                        "images": {
                                            "type": "array",
                                            "items": {"type": "string", "format": "binary"},
                                        },
                                    },
                                }
                            }
                        },
                    },
                    "responses": {
                        "200": {"description": "Report generated"},
                        "400": {"description": "Missing input"},
                    },
                }
            },
            "/api/get-instagram-data": {
                "get": {
                    "summary": "Fetch Instagram profile data",
                    "parameters": [
                        {"name": "username", "in": "query", "schema": {"type": "string"}},
                        {"name": "url", "in": "query", "schema": {"type": "string"}},
                    ],
                    "responses": {
                        "200": {"description": "Instagram profile fetched"},
                        "400": {"description": "Missing username/url"},
                    },
                },
                "post": {
                    "summary": "Fetch Instagram profile data",
                    "requestBody": {
                        "content": {
                            "application/json": {
                                "schema": {
                                    "type": "object",
                                    "properties": {
                                        "username": {"type": "string"},
                                        "url": {"type": "string"},
                                    },
                                }
                            }
                        }
                    },
                    "responses": {
                        "200": {"description": "Instagram profile fetched"},
                    },
                },
            },
            "/api/proxy-media": {
                "get": {
                    "summary": "Proxy remote media",
                    "parameters": [
                        {"name": "url", "in": "query", "required": True, "schema": {"type": "string"}}
                    ],
                    "responses": {
                        "200": {"description": "Media streamed"},
                        "400": {"description": "Invalid url"},
                    },
                }
            },
            "/api/template-preview": {
                "get": {
                    "summary": "Get template preview slides",
                    "responses": {"200": {"description": "Template preview metadata"}},
                }
            },
            "/api/report-preview/{filename}": {
                "get": {
                    "summary": "Get generated report preview slides",
                    "parameters": [
                        {"name": "filename", "in": "path", "required": True, "schema": {"type": "string"}}
                    ],
                    "responses": {"200": {"description": "Report preview metadata"}},
                }
            },
            "/api/download/{filename}": {
                "get": {
                    "summary": "Download generated PPTX",
                    "parameters": [
                        {"name": "filename", "in": "path", "required": True, "schema": {"type": "string"}}
                    ],
                    "responses": {"200": {"description": "PPTX file download"}},
                }
            },
        },
    }
    return jsonify(spec)


@app.route('/api/docs', methods=['GET'])
def swagger_docs():
    """Serve a lightweight Swagger UI page without extra backend dependencies."""
    server_url = request.host_url.rstrip('/')
    return Response(
        f"""<!doctype html>
<html>
  <head>
    <meta charset="utf-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1" />
    <title>Phyo Report Generator API Docs</title>
    <link rel="stylesheet" href="https://unpkg.com/swagger-ui-dist@5/swagger-ui.css" />
    <style>
      body {{
        margin: 0;
        background: #f6faf8;
      }}
      .topbar {{
        display: none;
      }}
    </style>
  </head>
  <body>
    <div id="swagger-ui"></div>
    <script src="https://unpkg.com/swagger-ui-dist@5/swagger-ui-bundle.js"></script>
    <script>
      window.ui = SwaggerUIBundle({{
        url: "{server_url}/api/openapi.json",
        dom_id: "#swagger-ui",
        deepLinking: true,
        presets: [SwaggerUIBundle.presets.apis],
      }});
    </script>
  </body>
</html>""",
        mimetype="text/html",
    )


@app.route('/api/prompt-presets', methods=['GET'])
def get_prompt_presets():
    """Return the built-in OpenRouter prompt presets."""
    return jsonify({
        "success": True,
        "default_template": "master",
        "templates": OPENROUTER_PROMPT_TEMPLATES,
    })


@app.route('/api/build-prompt', methods=['POST'])
def build_prompt():
    """Build a formatted prompt from one of the OpenRouter prompt presets."""
    body = request.get_json(silent=True) or {}
    template_name = body.get("template", "master")
    data = body.get("data", {}) or {}

    try:
        prompt = build_openrouter_prompt(template_name, data)
    except ValueError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400

    return jsonify({
        "success": True,
        "template": template_name,
        "prompt": prompt,
    })


@app.route('/api/generate-report', methods=['POST'])
def generate_report():
    """Main endpoint to generate campaign report"""
    try:
        warnings = []
        # Get form data
        prompt = request.form.get('prompt', '')
        raw_budget_values = request.form.get('budget_inr_values', '').strip()
        raw_images = request.files.getlist('images')
        images = [make_in_memory_upload(image) for image in raw_images]
        instagram_post_url = request.form.get('instagram_post_url', '').strip()
        instagram_profile_url = request.form.get('instagram_profile_url', '').strip()
        instagram_username = request.form.get('instagram_username', '').strip()
        youtube_post_url   = request.form.get('youtube_post_url',  '').strip()
        content_urls = _extract_urls_from_text(instagram_post_url, dedupe=False)
        max_content_urls = 10
        all_instagram_post_urls = [
            url for url in content_urls
            if "instagram.com" in url.lower() and any(token in url.lower() for token in ["/reel/", "/p/", "/tv/"])
        ]
        all_youtube_urls = [url for url in content_urls if "youtube.com" in url.lower() or "youtu.be" in url.lower()]
        instagram_post_urls_raw = all_instagram_post_urls[:max_content_urls]
        youtube_urls_raw = all_youtube_urls[:max_content_urls]
        instagram_post_urls = list(dict.fromkeys(instagram_post_urls_raw))
        youtube_urls = list(dict.fromkeys(youtube_urls_raw))

        budget_inr_values = []
        if raw_budget_values:
            try:
                parsed_budget_values = json.loads(raw_budget_values)
                if isinstance(parsed_budget_values, list):
                    budget_inr_values = [str(value).strip() if value is not None else "" for value in parsed_budget_values[:max_content_urls]]
            except Exception:
                quoted_values = re.findall(r'"([^"]*)"', raw_budget_values)
                if quoted_values:
                    budget_inr_values = [value.strip() for value in quoted_values[:max_content_urls]]
                else:
                    fallback_values = [
                        value.strip().strip('"').strip("'")
                        for value in re.split(r'[\n,]+', raw_budget_values)
                        if value.strip().strip('"').strip("'")
                    ]
                    budget_inr_values = fallback_values[:max_content_urls]

        # Validate: need at least one input
        if not prompt and not images and not instagram_post_urls and not instagram_profile_url and not instagram_username and not youtube_post_url and not youtube_urls:
            return jsonify({"success": False, "error": "Provide a prompt, images, an Instagram username/profile URL, or a post URL"}), 400

        print(f"\n{'='*60}")
        print(f"[START] NEW REQUEST")
        print(f"{'='*60}")
        print(f"Prompt: {prompt}")
        if budget_inr_values: print(f"Budget INR values: {budget_inr_values}")
        print(f"Images: {len(images)} uploaded")
        if instagram_post_urls: print(f"Instagram URLs: {instagram_post_urls}")
        if instagram_profile_url: print(f"Instagram Profile URL: {instagram_profile_url}")
        if instagram_username: print(f"Instagram Username: {instagram_username}")
        if youtube_post_url or youtube_urls: print(f"YouTube URL(s):   {youtube_urls or [youtube_post_url]}")

        # Step 1: Extract metrics from text prompt
        print("[NOTE] Extracting metrics from text...")
        text_metrics = extract_metrics_from_text(prompt) if prompt else {}
        prompt_context = extract_prompt_context(prompt) if prompt else {}
        cleaned_budget_values = []
        for value in budget_inr_values:
            cleaned_value = str(value).replace(",", "").strip()
            cleaned_budget_values.append(cleaned_value)

        numeric_budget_values = []
        for value in cleaned_budget_values:
            try:
                numeric_budget_values.append(float(value)) if value else None
            except Exception:
                continue

        if numeric_budget_values:
            prompt_context.setdefault("financial", {})
            total_budget_sum = sum(numeric_budget_values)
            prompt_context["financial"]["totalBudget"] = str(int(total_budget_sum)) if float(total_budget_sum).is_integer() else f"{total_budget_sum:.2f}"
            prompt_context["financial"]["budgetCurrency"] = "INR"

        # Step 2: Run slow independent fetches in parallel.
        instagram_post = {}
        instagram_profile = {}
        youtube_metrics = {}
        ocr_metrics = {}
        extracted_data = create_default_data()
        if instagram_post_urls_raw:
            extracted_data["requestedInstagramPostUrls"] = list(instagram_post_urls_raw)
            extracted_data["processedInstagramPostUrls"] = list(instagram_post_urls)
        if youtube_urls:
            extracted_data["requestedYoutubeUrls"] = list(youtube_urls)
        if len(all_instagram_post_urls) > max_content_urls:
            warnings.append(
                f"You pasted {len(all_instagram_post_urls)} Instagram URLs. The report currently processes the first {max_content_urls} URLs for speed and cleaner combined reporting."
            )
        if len(all_youtube_urls) > max_content_urls:
            warnings.append(
                f"You pasted {len(all_youtube_urls)} YouTube URLs. The report currently processes the first {max_content_urls} URLs."
            )

        yt_url = youtube_post_url or (youtube_urls[0] if youtube_urls else None)
        if not yt_url and prompt and ("youtube.com" in prompt or "youtu.be" in prompt):
            yt_url = prompt
        profile_ref = instagram_profile_url or instagram_username
        profile_username = (
            profile_ref.rstrip("/").split("/")[-1]
            if profile_ref and "instagram.com" in profile_ref.lower()
            else profile_ref
        )

        future_map = {}
        max_workers = 1 + int(bool(instagram_post_urls)) + int(bool(profile_username)) + int(bool(HAS_YT_DLP and yt_url)) + int(bool(images)) + int(bool(images))
        with ThreadPoolExecutor(max_workers=max(1, min(6, max_workers))) as executor:
            if instagram_post_urls:
                print(f"[INSTAGRAM POSTS] Queueing Bright Data batch fetch for {len(instagram_post_urls)} URL(s)")
                future_map["instagram_posts"] = executor.submit(fetch_instagram_posts_batch, instagram_post_urls)

            if profile_username:
                print(f"[INSTAGRAM PROFILE] Fetching profile data for: {profile_ref}")
                future_map["instagram_profile"] = executor.submit(fetch_instagram_data, profile_username)

            if HAS_YT_DLP and yt_url:
                print("[VIDEO] Fetching YouTube metrics...")
                future_map["youtube_metrics"] = executor.submit(fetch_youtube_metrics, yt_url)

            if images:
                future_map["ocr_metrics"] = executor.submit(extract_metrics_from_images_ocr, clone_uploads(images))
                future_map["vision_extract"] = executor.submit(analyze_images_with_gpt4, clone_uploads(images), prompt)

            instagram_post_results = []
            if "instagram_posts" in future_map:
                instagram_batch = future_map["instagram_posts"].result() or {}
                instagram_post_results = instagram_batch.get("results", []) or []
                invalid_urls = instagram_batch.get("invalid_urls", []) or []
                for invalid_url in invalid_urls:
                    warnings.append(
                        f"Skipped invalid Instagram post URL: {invalid_url}"
                    )

                fetched_url_keys = set()
                for post_data in instagram_post_results:
                    status_payload = post_data.get("urlStatus", {}) or {}
                    status_url = (status_payload.get("url") or "").strip()
                    if status_url:
                        fetched_url_keys.add(status_url.lower())

                    if post_data.get("hasMetrics"):
                        print(
                            f"[OK] Post metrics fetched for {status_url or 'Instagram URL'} - "
                            f"likes: {post_data['instagram'].get('likes')}, "
                            f"views: {post_data['instagram'].get('views')}"
                        )
                    else:
                        print(f"[OK] Partial Instagram metadata fetched for {status_url or 'Instagram URL'}")

                for url in instagram_post_urls:
                    if (url or "").strip().lower() not in fetched_url_keys:
                        warnings.append(
                            f"Bright Data did not return usable Instagram metrics for {url}, so the backend skipped that URL and continued with the remaining sources."
                        )
                if instagram_post_urls and not instagram_post_results:
                    warnings.append(
                        "Bright Data returned no matched post/reel rows. "
                        "Set correct dataset IDs: BRIGHTDATA_INSTAGRAM_REELS_DATASET_ID for /reel/ URLs and BRIGHTDATA_INSTAGRAM_POSTS_DATASET_ID for /p/ URLs."
                    )

            if len(instagram_post_results) > 1:
                instagram_post = merge_instagram_post_results(instagram_post_results) or {}
            elif len(instagram_post_results) == 1:
                instagram_post = instagram_post_results[0]

            if "instagram_profile" in future_map:
                instagram_profile = future_map["instagram_profile"].result()
                if not instagram_profile:
                    warnings.append(
                        "Bright Data did not return usable Instagram profile data, so the backend skipped profile metrics."
                    )

            if "youtube_metrics" in future_map:
                youtube_metrics = future_map["youtube_metrics"].result() or {}
                if not youtube_metrics:
                    warnings.append(
                        "YouTube URL fetch did not return usable metrics, so the backend skipped it and continued with the remaining sources."
                    )

            if "ocr_metrics" in future_map:
                ocr_metrics = future_map["ocr_metrics"].result() or {}

            if "vision_extract" in future_map:
                extracted_data = future_map["vision_extract"].result() or create_default_data()

        if instagram_post_urls_raw:
            extracted_data["requestedInstagramPostUrls"] = list(instagram_post_urls_raw)
            extracted_data["processedInstagramPostUrls"] = list(instagram_post_urls)
        if youtube_urls:
            extracted_data["requestedYoutubeUrls"] = list(youtube_urls)

        if images and not has_meaningful_metrics(extracted_data) and not ocr_metrics:
            warnings.append(
                "The uploaded images did not expose readable metrics, so the backend ignored them for metrics and used only prompt text and other available sources."
            )

        if not images and instagram_post:
            remote_media_uploads = build_remote_media_uploads(instagram_post)
            if remote_media_uploads:
                print(f"[IMAGE] Using {len(remote_media_uploads)} Bright Data media image(s) as fallback extraction input")
                remote_ocr_metrics = extract_metrics_from_images_ocr(clone_uploads(remote_media_uploads))
                if remote_ocr_metrics:
                    for key, value in remote_ocr_metrics.items():
                        if value and value != "N/A" and not ocr_metrics.get(key):
                            ocr_metrics[key] = value
                if not has_meaningful_metrics(extracted_data):
                    try:
                        extracted_data = analyze_images_with_gpt4(clone_uploads(remote_media_uploads), prompt) or extracted_data
                    except Exception as remote_image_error:
                        print(f"[WARNING] Remote image fallback analysis failed: {remote_image_error}")

        # Treat the user's prompt as the source of truth for campaign identity fields.
        extracted_data = apply_prompt_context(extracted_data, prompt_context, force_identity=True)

        # Step 4b: Fill top-level fields from prompt text if Vision AI left them empty
        import re as _re2

        def _extract_from_prompt(patterns, text):
            for p in patterns:
                m = _re2.search(p, text, _re2.IGNORECASE)
                if m:
                    return m.group(1).strip()
            return ""

        if prompt:
            if not extracted_data.get("campaignName") or extracted_data["campaignName"] in ("", "Campaign Report"):
                val = _extract_from_prompt([
                    r'report\s+for\s+([^\n,\.]+)',
                    r'campaign(?:\s+name)?[:\s]+([^\n,\.]+)',
                ], prompt)
                if val:
                    extracted_data["campaignName"] = val.strip()[:60]

            # Only extract brand if explicitly stated with "brand:" prefix — never guess from campaign name
            if not extracted_data.get("brand") or extracted_data["brand"] in ("", "Unknown Brand"):
                val = _extract_from_prompt([r'\bbrand[:\s]+([^\n,\.]+)'], prompt)
                if val:
                    extracted_data["brand"] = val.strip()[:40]

            if not extracted_data.get("creator"):
                val = _extract_from_prompt([
                    r'(?:influencer|creator)\s*[:\s]\s*([^\n,\.]{2,})',          # "influencer: Name"
                    r'with\s+((?:influencer|creator)\s+[^\n,\.]+?)(?:\s*,|\s+budget|\s+ran|$)',  # "with Influencer X"
                ], prompt)
                # Reject single-letter or very short matches (e.g. "X")
                if val and len(val.strip()) > 1:
                    extracted_data["creator"] = val.strip()[:40]

            # Budget → financial
            if not extracted_data.get("financial", {}).get("totalBudget"):
                bval = _extract_from_prompt([r'budget\s*\$?([\d,]+)'], prompt)
                if bval:
                    extracted_data.setdefault("financial", {})["totalBudget"] = bval.replace(",", "")

            # Dates
            if not extracted_data.get("startDate") or not extracted_data.get("endDate"):
                # Match "June 1-30" or "June 1 - June 30" or "June 1 to June 30"
                range_match = _re2.search(
                    r'ran\s+([A-Za-z]+)\s+(\d+)\s*[-–to]+\s*(?:([A-Za-z]+)\s+)?(\d+)',
                    prompt, _re2.IGNORECASE
                )
                if range_match:
                    month_start = range_match.group(1)
                    day_start   = range_match.group(2)
                    month_end   = range_match.group(3) or month_start
                    day_end     = range_match.group(4)
                    if not extracted_data.get("startDate"):
                        extracted_data["startDate"] = f"{month_start} {day_start}"
                    if not extracted_data.get("endDate"):
                        extracted_data["endDate"] = f"{month_end} {day_end}"
                else:
                    if not extracted_data.get("startDate"):
                        val = _extract_from_prompt([r'ran\s+([A-Za-z]+\s+\d+)'], prompt)
                        if val:
                            extracted_data["startDate"] = val

        extracted_data = apply_prompt_context(extracted_data, prompt_context, force_identity=True)

        # Step 5: Merge all metrics from prompt, images, and URLs.
        # Priority (highest → lowest):
        #   1. Bright Data / direct URL metrics with non-zero values
        #   2. Vision / OCR image metrics
        #   3. Prompt text metrics
        #   4. YouTube API fills YouTube fields
        print("[MERGE] Merging metrics from prompt, images, and URLs...")

        extracted_data.setdefault("instagram", {})
        extracted_data.setdefault("youtube", {})

        def fill_empty(target_dict, source_dict):
            """Copy from source into target only where target value is empty/missing."""
            for key, value in source_dict.items():
                if _has_metric_value(value) and not _has_metric_value(target_dict.get(key)):
                    target_dict[key] = value

        def fill_prefer_source(target_dict, source_dict):
            """Copy source into target when the source value is meaningful."""
            for key, value in source_dict.items():
                if _has_metric_value(value):
                    target_dict[key] = value

        def _has_metric_value(value):
            if value is None:
                return False
            if isinstance(value, str):
                return value.strip() not in ("", "0", "0.0", "0%")
            return bool(value)

        prompt_instagram = {}
        prompt_to_ig = {
            "views": "views",
            "likes": "likes",
            "shares": "shares",
            "saves": "saves",
            "engagement_rate": "engagementRate",
        }
        for text_key, ig_key in prompt_to_ig.items():
            value = text_metrics.get(text_key)
            if _has_metric_value(value):
                prompt_instagram[ig_key] = value

        image_instagram = {}
        ocr_to_ig = {
            "views": "views",
            "likes": "likes",
            "shares": "shares",
            "saves": "saves",
        }
        for ocr_key, ig_key in ocr_to_ig.items():
            value = ocr_metrics.get(ocr_key)
            if _has_metric_value(value):
                image_instagram[ig_key] = value

        # Keep useful metrics found from screenshots / OCR / prompt.
        fill_empty(extracted_data["instagram"], prompt_instagram)
        fill_empty(extracted_data["instagram"], image_instagram)

        # Keep financial totals from prompt/OCR text.
        if _has_metric_value(text_metrics.get("total_engagement")) and not _has_metric_value(extracted_data.get("performance", {}).get("totalEngagement")):
            extracted_data.setdefault("performance", {})["totalEngagement"] = text_metrics.get("total_engagement")
        if _has_metric_value(ocr_metrics.get("total_engagement")) and not _has_metric_value(extracted_data.get("performance", {}).get("totalEngagement")):
            extracted_data.setdefault("performance", {})["totalEngagement"] = ocr_metrics.get("total_engagement")

        # Financial details from text go into financial
        fin_key_map = {
            "budget": "totalBudget",
            "budget_currency": "budgetCurrency",
            "cpv": "cpv",
            "cpe": "cpe",
            "cpc": "cpc",
            "cpc_goal": "cpcGoal",
            "cpv_goal": "cpvGoal",
            "cpc_calculation": "cpcCalculation",
            "cpv_calculation": "cpvCalculation",
        }
        for txt_key, fin_key in fin_key_map.items():
            val = text_metrics.get(txt_key) or ocr_metrics.get(txt_key)
            if val and val != "N/A" and not extracted_data.get("financial", {}).get(fin_key):
                extracted_data.setdefault("financial", {})[fin_key] = val

        clicks_value = text_metrics.get("clicks") or ocr_metrics.get("clicks")
        if _has_metric_value(clicks_value) and not _has_metric_value(extracted_data.get("performance", {}).get("totalClicks")):
            extracted_data.setdefault("performance", {})["totalClicks"] = clicks_value

        # Instagram post URL — exact metrics for this specific post/reel.
        if instagram_post:
            ig_post_metrics = instagram_post.get("instagram", {})
            fill_prefer_source(extracted_data["instagram"], ig_post_metrics)
            creator_candidates = instagram_post.get("creatorNames") or []
            if not extracted_data.get("creator"):
                extracted_data["creator"] = (
                    " + ".join(creator_candidates[:3])
                    or instagram_post.get("fullName")
                    or instagram_post.get("username", "")
                )
            if (not extracted_data.get("brand") or extracted_data.get("brand") == "Unknown Brand") and instagram_post.get("brandName"):
                extracted_data["brand"] = instagram_post.get("brandName")
            if not extracted_data.get("deliverables"):
                extracted_data["deliverables"] = instagram_post.get("mediaType", "")
            if instagram_post.get("postImage"):
                extracted_data["postImage"] = instagram_post.get("postImage")
            if instagram_post.get("videoUrl"):
                extracted_data["videoUrl"] = instagram_post.get("videoUrl")
            if instagram_post.get("brandLogo"):
                extracted_data["brandLogo"] = instagram_post.get("brandLogo")
            if instagram_post.get("postUrls"):
                extracted_data["instagramPostUrls"] = instagram_post.get("postUrls")
            if instagram_post.get("urlResults"):
                extracted_data["instagramUrlResults"] = instagram_post.get("urlResults")
            if instagram_post.get("summary"):
                extracted_data["instagramCombinedSummary"] = instagram_post.get("summary")
            extracted_data["instagramSource"] = instagram_post.get("source", "instagram_url")
            if cleaned_budget_values:
                budget_per_url_map = {}
                for idx, url in enumerate(instagram_post_urls_raw):
                    normalized_url = (url or "").strip().lower()
                    if not normalized_url:
                        continue
                    budget_value_raw = cleaned_budget_values[idx] if idx < len(cleaned_budget_values) else ""
                    try:
                        budget_value = float(str(budget_value_raw).replace(",", "").strip()) if budget_value_raw else 0.0
                    except Exception:
                        budget_value = 0.0
                    if budget_value <= 0:
                        continue
                    budget_per_url_map.setdefault(normalized_url, {"url": url, "total": 0.0})
                    budget_per_url_map[normalized_url]["total"] += budget_value

                extracted_data["budgetPerUrl"] = [
                    {
                        "url": payload["url"],
                        "budget": str(int(payload["total"])) if float(payload["total"]).is_integer() else f"{payload['total']:.2f}",
                        "currency": "INR",
                    }
                    for payload in budget_per_url_map.values()
                ]
            if (not extracted_data.get("campaignName") or extracted_data.get("campaignName") == "Campaign Report"):
                campaign_candidate = (
                    instagram_post.get("postTitle")
                    or instagram_post.get("captionText", "").splitlines()[0]
                    or ""
                ).strip()
                if campaign_candidate:
                    extracted_data["campaignName"] = campaign_candidate[:80]

        # Instagram profile dataset — fill remaining Instagram fields only.
        if instagram_profile:
            fill_empty(extracted_data["instagram"], instagram_profile.get("instagram", {}))
            if not extracted_data.get("creator"):
                extracted_data["creator"] = (
                    instagram_profile.get("fullName") or instagram_profile.get("username", "")
                )
            extracted_data["instagramSource"] = extracted_data.get("instagramSource") or instagram_profile.get("source", "brightdata_profile_dataset")

        # YouTube API — fill empty youtube fields
        youtube_schema_keys = {"views", "likes", "comments", "shares", "watchTime", "ctr"}
        for key, value in youtube_metrics.items():
            if not value or value == "N/A":
                continue
            if key == "_title":
                if not extracted_data.get("campaignName"):
                    extracted_data["campaignName"] = value
            elif key == "_creator_name":
                if not extracted_data.get("creator"):
                    extracted_data["creator"] = value
            elif key in youtube_schema_keys:
                if not extracted_data["youtube"].get(key):
                    extracted_data["youtube"][key] = value
        
        # Step 5b: Propagate instagram → performance totals and creators
        ig = extracted_data.get("instagram", {})
        perf = extracted_data.setdefault("performance", {})

        # Map final instagram fields → performance totals
        ig_to_perf = {
            "views":          "totalViews",
            "likes":          "totalLikes",
            "comments":       "totalComments",
            "shares":         "totalShares",
            "saves":          "totalSaves",
            "reach":          "totalReach",
        }
        for ig_key, perf_key in ig_to_perf.items():
            if _has_metric_value(ig.get(ig_key)):
                perf[perf_key] = ig[ig_key]
        combined_summary = extracted_data.get("instagramCombinedSummary", {}) or {}
        combined_totals = combined_summary.get("totals", {}) or {}
        if combined_totals:
            total_perf_map = {
                "views": "totalViews",
                "likes": "totalLikes",
                "comments": "totalComments",
                "shares": "totalShares",
                "saves": "totalSaves",
                "reach": "totalReach",
            }
            for total_key, perf_key in total_perf_map.items():
                if _has_metric_value(combined_totals.get(total_key)):
                    perf[perf_key] = combined_totals[total_key]

        # Calculate totalInteractions = likes + comments + shares + saves
        def _safe_int(v):
            """Return int value only if v is a non-empty, non-zero string."""
            if not v or str(v).strip() in ("", "0"):
                return 0
            try:
                return int(float(str(v).replace('%', '').replace(',', '')))
            except Exception:
                return 0

        if not perf.get("totalInteractions"):
            total_interact = sum(_safe_int(ig.get(k)) for k in ["likes", "comments", "shares", "saves"])
            if total_interact > 0:
                perf["totalInteractions"] = str(total_interact)

        def _safe_float(v):
            if v in (None, "", "N/A"):
                return 0.0
            try:
                return float(str(v).replace('%', '').replace(',', '').strip())
            except Exception:
                return 0.0

        # Formula-based financial calculations
        financial = extracted_data.setdefault("financial", {})
        budget_val = _safe_float(financial.get("totalBudget"))
        views_val = _safe_float(perf.get("totalViews") or ig.get("views"))
        interactions_val = _safe_float(perf.get("totalInteractions"))
        clicks_val = _safe_float(perf.get("totalClicks"))
        if not financial.get("budgetCurrency"):
            financial["budgetCurrency"] = "INR"

        # Derive clicks from YouTube CTR and views when explicit clicks are missing.
        if clicks_val <= 0:
            yt_views = _safe_float(extracted_data.get("youtube", {}).get("views"))
            yt_ctr = _safe_float(extracted_data.get("youtube", {}).get("ctr"))
            if yt_views > 0 and yt_ctr > 0:
                clicks_val = round(yt_views * yt_ctr / 100.0, 2)
                perf["totalClicks"] = str(int(clicks_val)) if clicks_val.is_integer() else str(clicks_val)

        def _format_formula_metric(value):
            rounded = round(value, 2)
            return str(int(rounded)) if float(rounded).is_integer() else f"{rounded:.2f}"

        # Always prefer exact formula-derived numeric values in final output.
        if budget_val > 0 and views_val > 0:
            financial["cpv"] = _format_formula_metric(budget_val / views_val)

        if budget_val > 0 and interactions_val > 0:
            financial["cpe"] = _format_formula_metric(budget_val / interactions_val)

        if budget_val > 0 and clicks_val > 0:
            financial["cpc"] = _format_formula_metric(budget_val / clicks_val)

        # Propagate into creators list — dedupe and sync final instagram values
        creator_name = extracted_data.get("creator", "")
        creators = extracted_data.get("creators", [])
        combined_creator_names = (instagram_post.get("creatorNames") if isinstance(instagram_post, dict) else []) or []
        combined_creator_entries = (instagram_post.get("creatorEntries") if isinstance(instagram_post, dict) else []) or []
        if combined_creator_entries:
            creators = [dict(entry) for entry in combined_creator_entries]
            extracted_data["creators"] = creators
        elif not creators and combined_creator_names:
            creators = [{"name": name} for name in combined_creator_names]
            extracted_data["creators"] = creators
        elif not creators and creator_name:
            creators = [{"name": creator_name}]
            extracted_data["creators"] = creators
        deduped_creators = []
        creator_index_by_key = {}
        budget_by_url = {}
        for idx, url in enumerate(instagram_post_urls_raw):
            normalized_url = (url or "").strip().lower()
            if not normalized_url:
                continue
            budget_raw = cleaned_budget_values[idx] if idx < len(cleaned_budget_values) else ""
            budget_value = _safe_float(budget_raw)
            if budget_value > 0:
                budget_by_url[normalized_url] = budget_by_url.get(normalized_url, 0.0) + budget_value

        for c in creators:
            if not isinstance(c, dict):
                continue
            normalized_creator = dict(c)
            name = (normalized_creator.get("name") or creator_name or "Creator").strip()
            normalized_creator["name"] = name
            for ig_key in ["views", "likes", "comments", "shares", "saves", "reach", "engagementRate"]:
                existing_val = normalized_creator.get(ig_key, "")
                ig_val = ig.get(ig_key, "")
                if not _has_metric_value(existing_val) and _has_metric_value(ig_val):
                    normalized_creator[ig_key] = ig_val

            creator_post_url = (normalized_creator.get("postUrl") or "").strip().lower()
            creator_budget_value = normalized_creator.get("budget", "")
            if not _has_metric_value(creator_budget_value) and creator_post_url in budget_by_url:
                normalized_creator["budget"] = _format_formula_metric(budget_by_url[creator_post_url])

            total = sum(_safe_int(normalized_creator.get(k)) for k in ["likes", "comments", "shares", "saves"])
            normalized_creator["interactions"] = str(total) if total > 0 else normalized_creator.get("interactions", "")
            normalized_creator["platform"] = "Instagram"

            creator_budget_numeric = _safe_float(normalized_creator.get("budget"))
            creator_views_numeric = _safe_float(normalized_creator.get("views"))
            creator_interactions_numeric = _safe_float(normalized_creator.get("interactions"))
            creator_clicks_numeric = _safe_float(normalized_creator.get("clicks"))
            if creator_budget_numeric > 0 and creator_views_numeric > 0:
                normalized_creator["cpv"] = _format_formula_metric(creator_budget_numeric / creator_views_numeric)
            if creator_budget_numeric > 0 and creator_interactions_numeric > 0:
                normalized_creator["cpe"] = _format_formula_metric(creator_budget_numeric / creator_interactions_numeric)
            if creator_budget_numeric > 0 and creator_clicks_numeric > 0:
                normalized_creator["cpc"] = _format_formula_metric(creator_budget_numeric / creator_clicks_numeric)

            creator_name_key = re.sub(r"[^a-z0-9]+", "", name.lower())
            creator_url_key = (normalized_creator.get("postUrl") or "").strip().lower()
            dedupe_key = creator_url_key or creator_name_key or f"creator_{len(deduped_creators)}"

            if dedupe_key in creator_index_by_key:
                existing_creator = deduped_creators[creator_index_by_key[dedupe_key]]
                for merge_key, merge_value in normalized_creator.items():
                    if not _has_metric_value(merge_value):
                        continue
                    if not _has_metric_value(existing_creator.get(merge_key)):
                        existing_creator[merge_key] = merge_value
                continue

            creator_index_by_key[dedupe_key] = len(deduped_creators)
            deduped_creators.append(normalized_creator)
        extracted_data["creators"] = deduped_creators
        if deduped_creators:
            creator_names = []
            seen_creator_name_keys = set()
            for creator_item in deduped_creators:
                creator_item_name = (creator_item.get("name", "") or "").strip()
                if not creator_item_name:
                    continue
                creator_item_key = creator_item_name.lower()
                if creator_item_key in seen_creator_name_keys:
                    continue
                seen_creator_name_keys.add(creator_item_key)
                creator_names.append(creator_item_name)
            extracted_data["creator"] = " + ".join(creator_names)[:120]

        # Preserve request/debug URL fields in final payload even if earlier extractors replaced the object.
        if instagram_post_urls_raw:
            extracted_data["requestedInstagramPostUrls"] = list(instagram_post_urls_raw)
            extracted_data["processedInstagramPostUrls"] = list(instagram_post_urls)
        if youtube_urls:
            extracted_data["requestedYoutubeUrls"] = list(youtube_urls)

        # Step 6: Populate PowerPoint template
        output_path = populate_powerpoint(extracted_data, images)

        # Step 7: Upload to Google Drive (optional)
        google_slides_link = None
        try:
            google_slides_link = upload_to_google_drive(output_path, os.path.basename(output_path))
        except Exception as e:
            print(f"[WARNING]  Google Drive upload failed: {str(e)}")

        # Validate and clean extracted data
        extracted_data = validate_and_clean_data(extracted_data)

        # Debug: Print what's being returned
        print(f"[OK] Returning extracted data:")
        print(f"   - campaignName: {extracted_data.get('campaignName')}")
        print(f"   - brand: {extracted_data.get('brand')}")
        print(f"   - instagram views: {extracted_data.get('instagram', {}).get('views')}")
        print(f"   - postImage present: {bool(extracted_data.get('postImage'))}")
        print(f"   - videoUrl present: {bool(extracted_data.get('videoUrl'))}")
        print(f"   - creators: {len(extracted_data.get('creators', []))} creator(s)")
        print(f"   - Keys: {list(extracted_data.keys())}")
        # Build warnings list for the frontend
        ig = extracted_data.get("instagram", {})
        has_ig_metrics = any(ig.get(k) for k in ["views", "likes", "comments", "reach"])
        if not has_ig_metrics:
            warnings.append("No Instagram metrics found from Bright Data. Provide a valid Instagram profile URL/username or post/reel URL.")

        return jsonify({
            "success": True,
            "message": "Report generated successfully!",
            "filename": os.path.basename(output_path),
            "google_slides_link": google_slides_link,
            "extracted_data": extracted_data,
            "warnings": warnings,
        })
        
    except Exception as e:
        print(f"[ERROR] Error in generate_report: {str(e)}")
        import traceback
        traceback.print_exc()
        
        return jsonify({
            "success": False,
            "error": str(e)
        }), 500


def _safe_int(value):
    if value in (None, "", False):
        return 0
    if isinstance(value, bool):
        return 0
    if isinstance(value, (int, float)):
        return int(value)
    import re
    digits = re.sub(r"[^\d]", "", str(value))
    return int(digits) if digits else 0


def _dig(obj, *path):
    cur = obj
    for key in path:
        if isinstance(cur, dict):
            cur = cur.get(key)
        elif isinstance(cur, list) and isinstance(key, int):
            index = key if key >= 0 else len(cur) + key
            if 0 <= index < len(cur):
                cur = cur[index]
            else:
                return None
        else:
            return None
    return cur


def _first_value(record, *paths):
    for path in paths:
        if not isinstance(path, tuple):
            path = (path,)  
        value = _dig(record, *path)
        if value not in (None, ""):
            return value
    return None


def _find_first_media_url(obj, media_kind):
    """Recursively search nested payloads for a likely image/video URL."""
    if media_kind not in ("image", "video"):
        return ""

    preferred_key_hints = {
        "image": ("display", "image", "thumbnail", "cover", "poster", "src"),
        "video": ("video", "playback", "stream", "mp4"),
    }
    banned_key_hints = {
        "image": ("profile_pic", "avatar", "icon"),
        "video": (),
    }
    preferred_exts = {
        "image": (".jpg", ".jpeg", ".png", ".webp", ".gif"),
        "video": (".mp4", ".mov", ".m3u8"),
    }

    def walk(value, parent_key=""):
        if isinstance(value, dict):
            for key, child in value.items():
                found = walk(child, str(key).lower())
                if found:
                    return found
        elif isinstance(value, list):
            for child in value:
                found = walk(child, parent_key)
                if found:
                    return found
        elif isinstance(value, str):
            lower = value.lower().strip()
            if lower.startswith("http"):
                key_ok = any(h in parent_key for h in preferred_key_hints[media_kind]) if parent_key else False
                ext_ok = any(ext in lower for ext in preferred_exts[media_kind])
                banned = any(h in parent_key for h in banned_key_hints[media_kind])
                if not banned and (key_ok or ext_ok):
                    return value
        return ""

    return walk(obj)

  
def _extract_brightdata_records(payload):
    if isinstance(payload, list):
        return payload
    if isinstance(payload, dict):
        for key in ["data", "results", "items", "output"]:
            value = payload.get(key)
            if isinstance(value, list):
                return value
        if payload:
            return [payload]
    return []


def _parse_brightdata_response(resp):
    """
    Bright Data can return JSON or JSONL (application/jsonl).
    Normalize both response formats into Python objects.
    """
    try:
        return resp.json()
    except Exception:
        text = (resp.text or "").strip()
        if not text:
            return []
        parsed_rows = []
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            try:
                parsed_rows.append(json.loads(line))
            except Exception:
                continue
        if parsed_rows:
            return parsed_rows
        raise


def _parse_instagram_shortcode(post_url):
    cleaned = (post_url or "").strip()
    if not cleaned:
        return None

    match = re.search(
        r"(?:https?://)?(?:www\.)?instagram\.com/(?:p|reel|tv)/([A-Za-z0-9_-]+)",
        cleaned,
        re.IGNORECASE,
    )
    if not match:
        return None

    shortcode = match.group(1)
    lowered = cleaned.lower()
    if "/reel/" in lowered:
        media_path = "reel"
        media_type = "Instagram Reel"
    elif "/tv/" in lowered:
        media_path = "tv"
        media_type = "Instagram Video"
    else:
        media_path = "p"
        media_type = "Instagram Post"

    canonical_url = f"https://www.instagram.com/{media_path}/{shortcode}/"
    return {
        "original_url": cleaned,
        "shortcode": shortcode,
        "media_path": media_path,
        "media_type": media_type,
        "canonical_url": canonical_url,
    }


def _extract_record_shortcode(record):
    if not isinstance(record, dict):
        return ""

    shortcode = (
        _first_value(
            record,
            "shortcode",
            ("post", "shortcode"),
            ("node", "shortcode"),
            ("xdt_shortcode_media", "shortcode"),
        )
        or ""
    ).strip()
    if shortcode:
        return shortcode

    candidate_url = str(
        _first_value(record, "url", "post_url", ("post", "url"), ("node", "url"), ("xdt_shortcode_media", "url"))
        or ""
    ).strip()
    parsed = _parse_instagram_shortcode(candidate_url)
    return (parsed or {}).get("shortcode", "")


def _fetch_brightdata_instagram_post_records(shortcode_items):
    if not shortcode_items:
        return {}
    default_reels_dataset_id = "gd_lk5ns7kz21pck8jpis"
    default_posts_dataset_id = "gd_lk5ns7kz21pck8jpis"

    reels_dataset_id = (
        os.environ.get("BRIGHTDATA_INSTAGRAM_REELS_DATASET_ID", "").strip()
        or os.environ.get("BRIGHTDATA_INSTAGRAM_POST_DATASET_ID", "").strip()
        or default_reels_dataset_id
    )
    posts_dataset_id = (
        os.environ.get("BRIGHTDATA_INSTAGRAM_POSTS_DATASET_ID", "").strip()
        or os.environ.get("BRIGHTDATA_INSTAGRAM_POST_DATASET_ID", "").strip()
        or default_posts_dataset_id
    )

    groups = {
        "reel": {"dataset_id": reels_dataset_id, "urls": [], "shortcode_by_url": {}},
        "post": {"dataset_id": posts_dataset_id, "urls": [], "shortcode_by_url": {}},
    }
    seen_urls = set()

    for item in shortcode_items:
        canonical_url = (item.get("canonical_url") or "").strip()
        shortcode = (item.get("shortcode") or "").strip().lower()
        media_path = (item.get("media_path") or "").strip().lower()
        if not canonical_url or not shortcode:
            continue
        normalized_url = canonical_url.rstrip("/").lower()
        if normalized_url in seen_urls:
            continue
        seen_urls.add(normalized_url)
        group_key = "reel" if media_path in {"reel", "tv"} else "post"
        groups[group_key]["urls"].append(canonical_url)
        groups[group_key]["shortcode_by_url"][normalized_url] = shortcode

    records_by_shortcode = {}
    total_requested = 0
    total_raw_records = 0

    for group_name, group_payload in groups.items():
        urls = group_payload["urls"]
        dataset_id = group_payload["dataset_id"]
        shortcode_by_url = group_payload["shortcode_by_url"]
        if not urls:
            continue

        total_requested += len(urls)
        try:
            records = _fetch_brightdata_records(urls, dataset_id)
        except Exception as e:
            print(f"[WARNING] Bright Data {group_name} dataset fetch failed (dataset={dataset_id}): {e}")
            continue

        total_raw_records += len(records)
        for record in records:
            if not isinstance(record, dict):
                continue

            media_node = record.get("xdt_shortcode_media")
            normalized_record = media_node if isinstance(media_node, dict) else record
            shortcode = _extract_record_shortcode(normalized_record) or _extract_record_shortcode(record)

            if not shortcode:
                record_url = str(
                    _first_value(normalized_record, "url", "post_url", ("post", "url"), ("node", "url"))
                    or _first_value(record, "url", "post_url", ("post", "url"), ("node", "url"))
                    or ""
                ).strip()
                parsed = _parse_instagram_shortcode(record_url)
                if parsed and parsed.get("shortcode"):
                    shortcode = parsed["shortcode"]
                elif record_url:
                    shortcode = shortcode_by_url.get(record_url.rstrip("/").lower(), "")

            if not shortcode:
                continue

            records_by_shortcode.setdefault(shortcode.lower(), []).append(normalized_record)

    print(
        f"[BRIGHTDATA] Post batch fetch complete | requested={total_requested} "
        f"| raw_records={total_raw_records} | mapped_shortcodes={len(records_by_shortcode)}"
    )

    if total_requested and total_raw_records and not records_by_shortcode:
        print(
            "[WARNING] Bright Data returned rows but none matched requested shortcodes. "
            "Likely wrong dataset_id for reels/posts."
        )
    return records_by_shortcode


def _extract_urls_from_text(value, dedupe=True):
    if not value:
        return []
    urls = re.findall(r"https?://[^\s,]+", value)
    if not dedupe:
        return [url.strip().rstrip(").,]") for url in urls if url.strip().rstrip(").,]")]
    deduped = []
    seen = set()
    for url in urls:
        cleaned = url.strip().rstrip(").,]")
        if cleaned and cleaned not in seen:
            seen.add(cleaned)
            deduped.append(cleaned)
    return deduped


def _clean_brand_candidate(value):
    if value is None:
        return ""
    if isinstance(value, (list, dict)):
        return ""
    cleaned = re.sub(r"\s+", " ", str(value)).strip().strip("@#")
    cleaned = re.sub(r"[^A-Za-z0-9&+ ._-]", " ", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip(" -_.")
    if not cleaned:
        return ""
    lowered = cleaned.lower()
    alias_map = {
        "flixbusindia": "FlixBus",
        "flixbus india": "FlixBus",
    }
    if lowered in alias_map:
        return alias_map[lowered]
    if lowered.endswith("india") and lowered.replace(" ", "") == "flixbusindia":
        return "FlixBus"
    if "flixbus" in lowered:
        return "FlixBus"
    if cleaned.islower() or cleaned.isupper():
        cleaned = " ".join(
            token.upper() if token.lower() in {"bbc", "tv"} else token.capitalize()
            for token in cleaned.split()
        )
    return cleaned[:40]


def _select_unique_brand_name(candidates):
    normalized = []
    seen = set()
    for candidate in candidates:
        cleaned = _clean_brand_candidate(candidate)
        if not cleaned:
            continue
        key = re.sub(r"[^a-z0-9]+", "", cleaned.lower())
        if not key or key in seen:
            continue
        seen.add(key)
        normalized.append(cleaned)
    if not normalized:
        return ""
    preferred = sorted(
        normalized,
        key=lambda item: (
            0 if any(ch.isupper() for ch in item) else 1,
            0 if " " not in item else 1,
            len(item),
        )
    )
    return preferred[0]


def merge_instagram_post_results(post_results):
    if not post_results:
        return None

    merged = {
        "mediaType": "",
        "source": "brightdata_shortcode_media_combined",
        "username": "",
        "fullName": "",
        "brandUsername": "",
        "brandName": "",
        "brandLogo": "",
        "creatorNames": [],
        "creatorEntries": [],
        "postImage": "",
        "videoUrl": "",
        "postUrls": [],
        "urlResults": [],
        "raw_data": {"brightdata_shortcode_media": []},
        "summary": {
            "postCount": 0,
            "totals": {},
            "averages": {},
        },
        "instagram": {
            "views": "",
            "likes": "",
            "comments": "",
            "shares": "",
            "saves": "",
            "reach": "",
            "impressions": "",
            "engagementRate": "",
        },
    }

    metric_keys = ["views", "likes", "comments", "shares", "saves", "impressions"]
    metric_totals = {key: 0 for key in metric_keys}
    metric_counts = {key: 0 for key in metric_keys}
    reach_values = []
    engagement_rates = []
    media_types = []
    creator_names = []
    creator_entries = []
    brand_candidates = []

    for item in post_results:
        if not isinstance(item, dict):
            continue

        ig = item.get("instagram", {}) or {}
        merged["postUrls"].extend(item.get("postUrls") or [])
        if item.get("urlStatus"):
            merged["urlResults"].append(item.get("urlStatus"))
        raw_item = item.get("raw_data", {}).get("brightdata_shortcode_media")
        merged["raw_data"]["brightdata_shortcode_media"].append(raw_item)

        if item.get("username") and not merged["username"]:
            merged["username"] = item["username"]
        if item.get("fullName") and not merged["fullName"]:
            merged["fullName"] = item["fullName"]
        if item.get("postImage") and not merged["postImage"]:
            merged["postImage"] = item["postImage"]
        if item.get("videoUrl") and not merged["videoUrl"]:
            merged["videoUrl"] = item["videoUrl"]
        if item.get("brandLogo") and not merged["brandLogo"]:
            merged["brandLogo"] = item["brandLogo"]
        if item.get("mediaType"):
            media_types.append(item["mediaType"])
        if item.get("fullName"):
            creator_names.append(item["fullName"])
        elif item.get("username"):
            creator_names.append(item["username"])
        creator_entry_name = item.get("fullName") or item.get("username") or ""
        if creator_entry_name:
            creator_entries.append({
                "name": creator_entry_name,
                "brand": item.get("brandName") or _infer_brand_from_post_record(raw_item) or "",
                "platform": "Instagram",
                "views": ig.get("views", ""),
                "likes": ig.get("likes", ""),
                "comments": ig.get("comments", ""),
                "shares": ig.get("shares", ""),
                "saves": ig.get("saves", ""),
                "reach": ig.get("reach", ""),
                "engagementRate": ig.get("engagementRate", ""),
                "interactions": str(sum(_safe_int(ig.get(k)) for k in ["likes", "comments", "shares", "saves"])) or "",
                "postUrl": (item.get("postUrls") or [""])[0],
                "postImage": item.get("postImage", ""),
                "videoUrl": item.get("videoUrl", ""),
            })
        if item.get("brandName"):
            brand_candidates.append(item.get("brandName"))
        if raw_item:
            inferred_brand = _infer_brand_from_post_record(raw_item)
            if inferred_brand:
                brand_candidates.append(inferred_brand)
            if not merged.get("brandLogo"):
                inferred_logo = _infer_brand_logo_url_from_post_record(raw_item)
                if inferred_logo:
                    merged["brandLogo"] = inferred_logo

        for key in metric_keys:
            value = _safe_int(ig.get(key))
            if value > 0:
                metric_totals[key] += value
                metric_counts[key] += 1

        reach_value = _safe_int(ig.get("reach"))
        if reach_value > 0:
            reach_values.append(reach_value)

        rate = ig.get("engagementRate")
        if rate not in (None, "", "0", "0%"):
            try:
                engagement_rates.append(float(str(rate).replace("%", "").strip()))
            except Exception:
                pass

    post_count = len([item for item in post_results if isinstance(item, dict)])
    merged["summary"]["postCount"] = post_count

    for key in metric_keys:
        total_value = metric_totals[key]
        average_value = round(total_value / metric_counts[key], 2) if metric_counts[key] else 0
        if total_value > 0:
            merged["summary"]["totals"][key] = str(total_value)
        if average_value > 0:
            merged["summary"]["averages"][key] = str(int(average_value)) if float(average_value).is_integer() else f"{average_value:.2f}"
            merged["instagram"][key] = merged["summary"]["averages"][key]

    if reach_values:
        avg_reach = round(sum(reach_values) / len(reach_values), 2)
        merged["summary"]["totals"]["reach"] = str(max(reach_values))
        merged["summary"]["averages"]["reach"] = str(int(avg_reach)) if float(avg_reach).is_integer() else f"{avg_reach:.2f}"
        merged["instagram"]["reach"] = merged["summary"]["averages"]["reach"]

    if engagement_rates:
        merged["instagram"]["engagementRate"] = f"{round(sum(engagement_rates) / len(engagement_rates), 2)}%"

    if media_types:
        merged["mediaType"] = " + ".join(sorted(set(media_types)))

    merged["brandName"] = _select_unique_brand_name(brand_candidates)

    merged["creatorNames"] = list(dict.fromkeys(name for name in creator_names if name))
    deduped_creator_entries = []
    seen_creator_entry_keys = set()
    for entry in creator_entries:
        creator_key = ((entry.get("name") or "").strip().lower(), (entry.get("postUrl") or "").strip().lower())
        if creator_key in seen_creator_entry_keys:
            continue
        seen_creator_entry_keys.add(creator_key)
        deduped_creator_entries.append(entry)
    merged["creatorEntries"] = deduped_creator_entries
    merged["postUrls"] = list(dict.fromkeys(merged["postUrls"]))
    has_metrics = any(merged["instagram"].get(k) for k in ["views", "likes", "comments", "shares", "saves", "reach", "impressions"])
    has_metadata = any([
        merged.get("brandName"),
        merged.get("fullName"),
        merged.get("username"),
        merged.get("postImage"),
        merged.get("videoUrl"),
        merged.get("creatorEntries"),
        merged.get("postUrls"),
    ])
    return merged if (has_metrics or has_metadata) else None


def _infer_brand_from_post_record(record):
    if not isinstance(record, dict):
        return ""
    candidates = [
        _first_value(record, "brand", "brand_name", ("brand", "name"), ("sponsor", "name"), ("sponsored_by", "name")),
        _first_value(record, ("coauthor_producers", 0, "username")),
        _first_value(record, ("coauthor_producers", 0, "full_name")),
        _first_value(record, "sponsor_user", "username"),
        _first_value(record, "sponsor_user", "full_name"),
        _first_value(record, ("edge_media_to_tagged_user", "edges", 0, "node", "user", "username")),
        _first_value(record, ("edge_media_to_tagged_user", "edges", 0, "node", "user", "full_name")),
    ]
    return _select_unique_brand_name(candidates)


def _infer_brand_logo_url_from_post_record(record):
    if not isinstance(record, dict):
        return ""
    candidates = [
        _first_value(record, "profile_image_link"),
        _first_value(record, "profile_pic_url"),
        _first_value(record, ("coauthor_producers", 0, "profile_pic_url")),
        _first_value(record, ("edge_media_to_tagged_user", "edges", 0, "node", "user", "profile_pic_url")),
        _first_value(record, "brand_logo_url"),
        _first_value(record, ("brand", "logo")),
        _first_value(record, ("sponsor_user", "profile_pic_url")),
    ]
    for candidate in candidates:
        if isinstance(candidate, str) and candidate.strip().startswith("http"):
            return candidate.strip()
    return ""


def _fetch_brightdata_records(urls, dataset_id, max_attempts=12, poll_delay=3):
    import requests as req
    import time

    brightdata_token = (
        os.environ.get("BRIGHTDATA_API_TOKEN", "").strip()
        or os.environ.get("BRIGHTDATA_API_KEY", "").strip()
    )
    if not brightdata_token:
        print("[WARNING] Bright Data token is not configured")
        return []

    if not dataset_id:
        print("[WARNING] Bright Data dataset id is not configured")
        return []

    resp = req.post(
        f"https://api.brightdata.com/datasets/v3/scrape?dataset_id={dataset_id}&notify=false&include_errors=true",
        headers={
            "Authorization": f"Bearer {brightdata_token}",
            "Content-Type": "application/json",
        },
        json={"input": [{"url": url} for url in urls]},
        timeout=90,
    )
    resp.raise_for_status()
    payload = _parse_brightdata_response(resp)
    records = _extract_brightdata_records(payload)
    if records:
        return records

    snapshot_id = ""
    if isinstance(payload, dict):
        snapshot_id = payload.get("snapshot_id") or payload.get("id") or ""
    if not snapshot_id:
        return []

    print(f"[BRIGHTDATA] Snapshot created: {snapshot_id}")
    for attempt in range(max_attempts):
        progress_resp = req.get(
            f"https://api.brightdata.com/datasets/v3/progress/{snapshot_id}",
            headers={"Authorization": f"Bearer {brightdata_token}"},
            timeout=30,
        )
        progress_resp.raise_for_status()
        progress_data = _parse_brightdata_response(progress_resp)
        status = progress_data.get("status", "")
        print(f"[BRIGHTDATA] Snapshot status: {status} (attempt {attempt + 1}/{max_attempts})")

        if status == "ready":
            snapshot_resp = req.get(
                f"https://api.brightdata.com/datasets/v3/snapshot/{snapshot_id}",
                headers={"Authorization": f"Bearer {brightdata_token}"},
                params={"format": "json"},
                timeout=60,
            )
            snapshot_resp.raise_for_status()
            return _extract_brightdata_records(_parse_brightdata_response(snapshot_resp))

        if status == "failed":
            print(f"[WARNING] Bright Data snapshot failed for {snapshot_id}")
            return []

        time.sleep(poll_delay)

    return []


def fetch_instagram_post_data(post_url, prefetched_records_by_shortcode=None):
    """Fetch metrics for a specific Instagram post/reel from Bright Data."""
    cache_key = ("instagram_post", (post_url or "").strip())
    cached = _cache_get(cache_key)
    if cached is not None:
        print(f"[CACHE] Using cached Instagram post data for: {post_url}")
        return cached

    parsed = _parse_instagram_shortcode(post_url)
    if not parsed:
        print(f"[ERROR] Cannot parse shortcode from: {post_url}")
        return None

    shortcode = parsed["shortcode"]
    canonical_url = parsed["canonical_url"]
    media_type = parsed["media_type"]
    print(f"[INSTAGRAM] Parsed shortcode {shortcode} from URL: {post_url}")

    merged = {
        "mediaType": media_type,
        "source": "",
        "username": "",
        "fullName": "",
        "brandUsername": "",
        "brandName": "",
        "creatorNames": [],
        "postImage": "",
        "videoUrl": "",
        "postTitle": "",
        "captionText": "",
        "hasMetrics": False,
        "hasMetadata": False,
        "postUrls": [post_url],
        "raw_data": {},
        "urlStatus": {
            "url": post_url,
            "shortcode": shortcode,
            "status": "pending",
            "hasMetrics": False,
            "hasMetadata": False,
            "creator": "",
            "brand": "",
        },
        "instagram": {
            "views": "",
            "likes": "",
            "comments": "",
            "shares": "",
            "saves": "",
            "reach": "",
            "impressions": "",
            "engagementRate": "",
        },
    }

    def _fill_post_field(key, value):
        if value not in (None, "") and not merged.get(key):
            merged[key] = value

    def _fill_instagram_field(key, value):
        if value not in (None, "") and not merged["instagram"].get(key):
            merged["instagram"][key] = str(value)

    if prefetched_records_by_shortcode is not None:
        records = prefetched_records_by_shortcode.get(shortcode.lower(), [])
    else:
        records = _fetch_brightdata_instagram_post_records([parsed]).get(shortcode.lower(), [])

    for record in records:
        if not isinstance(record, dict):
            continue

        record_shortcode = _extract_record_shortcode(record) or ""
        record_url = (
            _first_value(record, "url", "post_url", ("post", "url"), ("node", "url"))
            or ""
        )
        if record_shortcode and record_shortcode.lower() != shortcode.lower():
            continue
        if record_url:
            parsed_record_url = _parse_instagram_shortcode(str(record_url))
            if parsed_record_url and parsed_record_url.get("shortcode", "").lower() != shortcode.lower():
                continue

        likes = _safe_int(_first_value(record, "like_count", "likes", ("metrics", "likes"), ("statistics", "likes"), ("post", "like_count"), ("node", "like_count"), ("edge_media_preview_like", "count")))
        comments = _safe_int(_first_value(record, "comment_count", "comments_count", "comments", "num_comments", ("metrics", "comments"), ("statistics", "comments"), ("post", "comment_count"), ("node", "comment_count"), ("edge_media_to_parent_comment", "count"), ("edge_media_preview_comment", "count")))
        views = _safe_int(_first_value(record, "video_play_count", "play_count", "video_view_count", "view_count", "views", ("metrics", "views"), ("statistics", "views"), ("post", "video_play_count"), ("post", "play_count"), ("post", "video_view_count"), ("node", "video_view_count")))
        shares = _safe_int(_first_value(record, "share_count", "shares", ("metrics", "shares"), ("statistics", "shares"), ("post", "share_count")))
        saves = _safe_int(_first_value(record, "save_count", "saves", ("metrics", "saves"), ("statistics", "saves"), ("post", "save_count")))
        reach = _safe_int(_first_value(record, "reach", "followers", ("owner", "edge_followed_by", "count"), ("user", "follower_count"), ("owner", "follower_count")))
        impressions = _safe_int(_first_value(record, "impressions", ("post", "impressions")))
        username = _first_value(record, "user_posted", ("owner", "username"), ("user", "username"), "username", ("author", "username"), ("post", "owner_username")) or ""
        full_name = _first_value(record, "profile_name", ("owner", "full_name"), ("user", "full_name"), "full_name", ("author", "full_name"), ("post", "owner_full_name")) or ""
        engagement_rate = _first_value(record, "engagement_rate", "avg_engagement", ("metrics", "engagement_rate"), ("post", "engagement_rate")) or ""
        post_title = _first_value(record, "title", "description", ("post", "title"), ("node", "title")) or ""
        caption_text = _first_value(
            record,
            "caption",
            "description",
            ("post", "caption"),
            ("node", "caption"),
            ("edge_media_to_caption", "edges", 0, "node", "text"),
        ) or ""
        post_image = _first_value(
            record,
            ("display_resources", -1, "src"),
            ("node", "display_resources", -1, "src"),
            "display_url",
            "displayUrl",
            "image_url",
            "thumbnail",
            "thumbnail_src",
            "thumbnail",
            "thumbnail_url",
            ("post", "display_url"),
            ("post", "thumbnail_src"),
            ("node", "display_url"),
            ("node", "thumbnail_src"),
            ("image_versions2", "candidates", 0, "url"),
            ("post", "image_versions2", "candidates", 0, "url"),
            ("display_resources", 0, "src"),
            ("node", "display_resources", 0, "src"),
            ("additional_candidates", "first_frame", "url"),
            ("post", "additional_candidates", "first_frame", "url"),
            ("carousel_media", 0, "image_versions2", "candidates", 0, "url"),
            ("photos", 0),
        ) or _find_first_media_url(record, "image") or ""
        video_url = _first_value(
            record,
            "video_url",
            "videoUrl",
            ("post", "video_url"),
            ("node", "video_url"),
            ("video_versions", 0, "url"),
            ("post", "video_versions", 0, "url"),
            ("carousel_media", 0, "video_versions", 0, "url"),
        ) or _find_first_media_url(record, "video") or ""

        _fill_instagram_field("views", views)
        _fill_instagram_field("likes", likes)
        _fill_instagram_field("comments", comments)
        _fill_instagram_field("shares", shares)
        _fill_instagram_field("saves", saves)
        _fill_instagram_field("reach", reach)
        _fill_instagram_field("impressions", impressions)
        if engagement_rate:
            _fill_instagram_field("engagementRate", engagement_rate)
        _fill_post_field("username", username)
        _fill_post_field("fullName", full_name)
        if full_name:
            merged["creatorNames"] = [full_name]
        elif username:
            merged["creatorNames"] = [username]
        if not merged.get("brandName"):
            merged["brandName"] = _infer_brand_from_post_record(record)
        if not merged.get("brandLogo"):
            merged["brandLogo"] = _infer_brand_logo_url_from_post_record(record)
        _fill_post_field("postImage", post_image)
        _fill_post_field("videoUrl", video_url)
        _fill_post_field("postTitle", post_title)
        _fill_post_field("captionText", caption_text)
        merged["source"] = merged["source"] or "brightdata_shortcode_media"
        merged["raw_data"]["brightdata_shortcode_media"] = record
        print(
            f"[OK] Bright Data post data found for shortcode {shortcode} | "
            f"creator={full_name or username or 'n/a'} | "
            f"likes={likes} | comments={comments} | views={views} | reach={reach}"
        )
        break

    if not merged["instagram"].get("engagementRate"):
        try:
            views = _safe_int(merged["instagram"].get("views"))
            likes = _safe_int(merged["instagram"].get("likes"))
            comments = _safe_int(merged["instagram"].get("comments"))
            if views > 0 and (likes or comments):
                merged["instagram"]["engagementRate"] = f"{round((likes + comments) / views * 100, 2)}%"
        except Exception:
            pass

    has_metrics = any(merged["instagram"].get(k) for k in ["views", "likes", "comments", "shares", "saves", "reach", "impressions"])
    has_metadata = any([
        merged.get("username"),
        merged.get("fullName"),
        merged.get("brandName"),
        merged.get("brandLogo"),
        merged.get("postImage"),
        merged.get("videoUrl"),
        merged.get("postTitle"),
        merged.get("captionText"),
    ])
    merged["hasMetrics"] = has_metrics
    merged["hasMetadata"] = has_metadata
    merged["urlStatus"] = {
        "url": post_url,
        "shortcode": shortcode,
        "status": "metrics" if has_metrics else ("metadata_only" if has_metadata else "failed"),
        "hasMetrics": has_metrics,
        "hasMetadata": has_metadata,
        "creator": merged.get("fullName") or merged.get("username") or "",
        "brand": _select_unique_brand_name([merged.get("brandName")]) or "",
    }
    print(f"[INSTAGRAM] URL status for {shortcode}: {merged['urlStatus']}")
    if not has_metrics and not has_metadata:
        print(f"[WARNING] No usable Instagram post metrics or metadata returned from Bright Data for {canonical_url}")
        return None

    _cache_set(cache_key, merged)
    return merged


def fetch_instagram_posts_batch(post_urls):
    parsed_items = []
    invalid_urls = []
    for url in post_urls:
        parsed = _parse_instagram_shortcode(url)
        if parsed:
            parsed_items.append(parsed)
        else:
            invalid_urls.append(url)

    if not parsed_items:
        return {"results": [], "invalid_urls": invalid_urls}

    prefetched_records_by_shortcode = _fetch_brightdata_instagram_post_records(parsed_items)
    results = []
    for parsed in parsed_items:
        item = fetch_instagram_post_data(
            parsed.get("original_url"),
            prefetched_records_by_shortcode=prefetched_records_by_shortcode,
        )
        if item:
            results.append(item)

    return {"results": results, "invalid_urls": invalid_urls}


def fetch_instagram_data(username):
    """Fetch Instagram profile data using only the Bright Data profile dataset."""
    username = (username or "").strip().lstrip("@")
    if not username:
        return None
    cache_key = ("instagram_profile", username.lower())
    cached = _cache_get(cache_key)
    if cached is not None:
        print(f"[CACHE] Using cached Instagram profile data for: @{username}")
        return cached

    dataset_id = (
        os.environ.get("BRIGHTDATA_INSTAGRAM_PROFILE_DATASET_ID", "").strip()
        or os.environ.get("BRIGHTDATA_INSTAGRAM_DATASET_ID", "").strip()
        or "gd_l1vikfch901nx3by4"
    )
    profile_url = f"https://www.instagram.com/{username}/"

    try:
        records = _fetch_brightdata_records([profile_url], dataset_id)
    except Exception as e:
        print(f"[WARNING] Bright Data dataset fetch failed: {e}")
        return None

    for record in records:
        if not isinstance(record, dict):
            continue

        record_username = (_first_value(record, "username", ("user", "username"), ("owner", "username")) or "").strip().lstrip("@")
        record_url = str(_first_value(record, "url", "profile_url", ("user", "url")) or "")
        if record_username and record_username.lower() != username.lower():
            continue
        if record_url and f"/{username.lower()}" not in record_url.lower():
            continue

        followers = _safe_int(_first_value(record, "followers", "followers_count", ("user", "followers_count"), ("owner", "edge_followed_by", "count")))
        following = _safe_int(_first_value(record, "following", "following_count", ("user", "following_count"), ("owner", "edge_follow", "count")))
        post_count = _safe_int(_first_value(record, "posts_count", "posts", ("user", "media_count")))
        full_name = _first_value(record, "profile_name", "full_name", ("user", "full_name")) or username
        bio = _first_value(record, "biography", "bio", ("user", "biography")) or ""
        is_verified = bool(_first_value(record, "is_verified", ("user", "is_verified")))
        profile_pic = _first_value(record, "profile_image_link", "profile_pic_url", ("user", "profile_pic_url")) or ""

        posts = _first_value(record, "posts") or []
        if not isinstance(posts, list):
            posts = []

        likes_list = [_safe_int(post.get("likes")) for post in posts if isinstance(post, dict) and post.get("likes") not in (None, "")]
        comments_list = [_safe_int(post.get("comments")) for post in posts if isinstance(post, dict) and post.get("comments") not in (None, "")]
        views_list = [_safe_int(post.get("views")) for post in posts if isinstance(post, dict) and post.get("views") not in (None, "")]

        def avg(values):
            return int(sum(values) / len(values)) if values else 0

        avg_likes = avg(likes_list)
        avg_comments = avg(comments_list)
        avg_views = avg(views_list)
        engagement_rate = ""
        if followers > 0 and (avg_likes or avg_comments):
            engagement_rate = f"{round((avg_likes + avg_comments) / followers * 100, 2)}%"

        print(f"[OK] Bright Data profile data found for @{username}")
        result = {
            "source": "brightdata_profile_dataset",
            "username": record_username or username,
            "fullName": full_name,
            "bio": bio,
            "isVerified": is_verified,
            "profilePic": profile_pic,
            "followersCount": followers,
            "followingCount": following,
            "postCount": post_count,
            "instagram": {
                "views": str(avg_views) if avg_views else "",
                "likes": str(avg_likes) if avg_likes else "",
                "comments": str(avg_comments) if avg_comments else "",
                "shares": "",
                "saves": "",
                "reach": str(followers) if followers else "",
                "impressions": "",
                "engagementRate": engagement_rate,
            },
            "raw_data": record,
        }
        _cache_set(cache_key, result)
        return result

    print(f"[WARNING] Bright Data returned no usable Instagram profile data for @{username}")
    return None


@app.route('/api/get-instagram-data', methods=['GET', 'POST'])
def get_instagram_data():
    """
    GET  /api/get-instagram-data?username=cristiano
    GET  /api/get-instagram-data?url=https://www.instagram.com/cristiano/
    POST /api/get-instagram-data  body: {"username": "cristiano"} or {"url": "..."}
    Returns Instagram profile data from the Bright Data profile dataset.
    """
    if request.method == 'POST':
        body     = request.get_json(silent=True) or {}
        username = body.get('username', '').strip()
        profile_url = body.get('url', '').strip()
    else:
        username = request.args.get('username', '').strip()
        profile_url = request.args.get('url', '').strip()

    if not username and not profile_url:
        return jsonify({"success": False, "error": "username or url is required"}), 400

    if profile_url and "instagram.com" in profile_url.lower():
        username = profile_url.rstrip("/").split("/")[-1]

    username = username.lstrip('@')

    print(f"[INSTAGRAM] Fetching Bright Data profile data for @{username} ...")
    data = fetch_instagram_data(username)

    if data is None:
        return jsonify({
            "success": False,
            "error":   (
                "Could not fetch Instagram data from Bright Data. "
                "The account may be private, missing, or unavailable from the dataset."
            )
        }), 502

    return jsonify({"success": True, "data": data})


@app.route('/api/proxy-media')
def proxy_media():
    media_url = request.args.get('url', '').strip()
    if not media_url or not media_url.startswith(('http://', 'https://')):
        return jsonify({"success": False, "error": "valid url is required"}), 400

    try:
        import requests
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/123.0.0.0 Safari/537.36"
            ),
            "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
            "Referer": "https://www.instagram.com/",
        }
        resp = requests.get(media_url, headers=headers, timeout=30)
        resp.raise_for_status()
        content_type = resp.headers.get("Content-Type", "application/octet-stream")
        return Response(resp.content, content_type=content_type)
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 502


@app.route('/api/template-preview', methods=['GET'])
def get_template_preview():
    """Return latest PNG previews for the local PPT template."""
    try:
        preview_data = export_template_preview_images()
        return jsonify({"success": True, **preview_data})
    except Exception as e:
        print(f"[ERROR] Template preview export failed: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/template-preview/image/<filename>', methods=['GET'])
def get_template_preview_image(filename):
    """Serve one exported template slide preview image."""
    try:
        safe_name = os.path.basename(filename)
        if safe_name != filename or not re.fullmatch(r"Slide\d+\.PNG", safe_name, re.IGNORECASE):
            return jsonify({"success": False, "error": "Invalid preview filename"}), 400

        image_path = os.path.join(PREVIEW_DIR, safe_name)
        if not os.path.exists(image_path):
            return jsonify({"success": False, "error": "Preview image not found"}), 404

        return send_file(image_path, mimetype='image/png')
    except Exception as e:
        print(f"[ERROR] Error serving preview image: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/report-preview/<filename>', methods=['GET'])
def get_report_preview(filename):
    """Return latest PNG previews for a generated PPTX report."""
    try:
        preview_data = export_report_preview_images(filename)
        return jsonify({"success": True, **preview_data})
    except Exception as e:
        print(f"[ERROR] Report preview export failed: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/report-preview/image/<report_key>/<filename>', methods=['GET'])
def get_report_preview_image(report_key, filename):
    """Serve one exported generated-report preview image."""
    try:
        safe_key = os.path.basename(report_key)
        safe_name = os.path.basename(filename)
        if (
            safe_key != report_key
            or safe_name != filename
            or not re.fullmatch(r"[A-Za-z0-9_.-]+", safe_key)
            or not re.fullmatch(r"Slide\d+\.PNG", safe_name, re.IGNORECASE)
        ):
            return jsonify({"success": False, "error": "Invalid preview path"}), 400

        image_path = os.path.join(REPORT_PREVIEW_ROOT, safe_key, safe_name)
        if not os.path.exists(image_path):
            return jsonify({"success": False, "error": "Preview image not found"}), 404

        return send_file(image_path, mimetype='image/png')
    except Exception as e:
        print(f"[ERROR] Error serving generated preview image: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/download/<filename>')
def download_report(filename):
    """Download generated PowerPoint report"""
    try:
        file_path = os.path.join(OUTPUT_DIR, filename)
        
        if not os.path.exists(file_path):
            return jsonify({
                "success": False,
                "error": "File not found"
            }), 404
        
        return send_file(
            file_path,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        print(f"[ERROR] Error downloading file: {str(e)}")
        return jsonify({
            "success": False,
            "error": str(e)
        }), 500


if __name__ == '__main__':
    # Check if template exists
    if not os.path.exists(TEMPLATE_PATH):
        print(f"[WARNING]  WARNING: Template file not found: {TEMPLATE_PATH}")
        print("Please place your template file in the same directory as app.py")
    
    # Check if API key is set
    if not os.environ.get("OPENROUTER_API_KEY") and not os.environ.get("ANTHROPIC_API_KEY"):
        print("[WARNING]  WARNING: No API key found!")
        print("Please set OPENROUTER_API_KEY or ANTHROPIC_API_KEY in .env")
    
    print("\n[START] Starting Flask server...")
    print(f"[DIR] Output directory: {OUTPUT_DIR}")
    print(f"[FILE] Template: {TEMPLATE_PATH}")
    
    app.run(debug=True, port=5000, host='0.0.0.0')
